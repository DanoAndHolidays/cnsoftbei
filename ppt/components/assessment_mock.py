"""
PPT 评估页 mockup：替代缺失的 Assessment 截图。

绘制：顶部 4 个统计卡（已完成题数/正确率/学习时长/连续天数）
     + 左侧 6 维能力雷达图（matplotlib）
     + 右侧学习建议时间线
     + 底部 4 个模块进度条
"""

import os
from pathlib import Path
from pptx.util import Pt, Emu, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import theme
from .layout import hex_to_rgb, add_textbox, add_rect
from .shapes import add_card, add_color_block
from .radar_chart import render_radar
from .flow_diagram import build_node


def render_assessment_mock(slide, left, top, width, height):
    """
    在指定区域绘制一张评估页 mockup。

    区域建议 ≥ Pt(900, 500)，即 900pt 宽 500pt 高。
    """
    # ---- 整体卡片背景
    add_card(slide, left, top, width, height, fill=theme.WHITE, border=theme.BORDER, border_width=1.0)

    # ---- 顶部 4 个统计卡
    stats = [
        ("已做题数", "126", "题", theme.PRIMARY),
        ("正确率",   "78", "%",  theme.SUCCESS),
        ("学习时长", "32", "小时", theme.WARNING),
        ("连续天数", "12", "天", "#722ED1"),
    ]
    stat_w = (width - Pt(60)) // 4
    stat_h = Pt(70)
    stat_top = top + Pt(20)
    stat_left_start = left + Pt(20)
    for i, (label, value, unit, color) in enumerate(stats):
        x = stat_left_start + i * (stat_w + Pt(10))
        # 卡片
        add_card(slide, x, stat_top, stat_w - Pt(10), stat_h, fill=theme.ACCENT_BG, border=None)
        # 左侧色条
        add_rect(slide, x, stat_top, Pt(4), stat_h, fill=color)
        # 数字
        add_textbox(
            slide, x + Pt(12), stat_top + Pt(8), stat_w - Pt(20), Pt(32),
            text=value, font_size=24, bold=True, color=theme.PRIMARY_DARK,
        )
        # 标签
        add_textbox(
            slide, x + Pt(12), stat_top + Pt(40), stat_w - Pt(20), Pt(20),
            text=f"{label} ({unit})", font_size=11, color=theme.TEXT_MUTED,
        )

    # ---- 左侧雷达图（PNG 渲染 → 插入）
    radar_path = os.path.join(os.path.dirname(__file__), "..", "assets", "_mock_radar.png")
    radar_path = os.path.abspath(radar_path)
    render_radar(
        radar_path,
        labels=["知识基础", "认知风格", "易错偏好", "学习节奏", "兴趣方向", "学习习惯"],
        values=[82, 70, 65, 88, 75, 80],
        title="",
        color=theme.PRIMARY,
    )
    radar_left = left + Pt(20)
    radar_top = stat_top + stat_h + Pt(20)
    radar_w = (width - Pt(60)) // 2
    radar_h = height - (radar_top - top) - Pt(20)
    slide.shapes.add_picture(radar_path, radar_left, radar_top, width=radar_w, height=radar_h)

    # ---- 右侧学习建议时间线
    sug_left = radar_left + radar_w + Pt(20)
    sug_top = radar_top
    sug_w = width - (sug_left - left) - Pt(20)
    sug_h = radar_h
    add_card(slide, sug_left, sug_top, sug_w, sug_h, fill=theme.WHITE, border=theme.BORDER)
    add_textbox(
        slide, sug_left + Pt(16), sug_top + Pt(12), sug_w - Pt(32), Pt(28),
        text="智能学习建议", font_size=15, bold=True, color=theme.PRIMARY_DARK,
    )
    suggestions = [
        ("知识基础", "基础扎实", theme.SUCCESS),
        ("易错偏好", "易错于装饰器，建议专项练习", theme.WARNING),
        ("学习节奏", "建议每天 45 分钟连续学习", theme.PRIMARY),
        ("兴趣方向", "倾向 Web 开发，可拓展前后端", theme.PRIMARY),
    ]
    for i, (tag, content, color) in enumerate(suggestions):
        item_top = sug_top + Pt(50) + i * Pt(40)
        # 左侧色圆点
        dot_size = Pt(10)
        add_rect(slide, sug_left + Pt(20), item_top + Pt(8), dot_size, dot_size, fill=color)
        # 标签
        add_textbox(
            slide, sug_left + Pt(36), item_top, Pt(60), Pt(24),
            text=tag, font_size=12, bold=True, color=color,
        )
        # 内容
        add_textbox(
            slide, sug_left + Pt(100), item_top, sug_w - Pt(120), Pt(24),
            text=content, font_size=12, color=theme.TEXT,
        )


def cleanup_mock_assets():
    """清理中间生成的 mockup PNG（每次生成 PPT 后调用）"""
    p = os.path.join(os.path.dirname(__file__), "..", "assets", "_mock_radar.png")
    p = os.path.abspath(p)
    if os.path.exists(p):
        try:
            os.remove(p)
        except OSError:
            pass