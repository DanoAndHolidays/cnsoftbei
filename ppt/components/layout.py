"""
PPT 页面框架 · 学术商务版（960×540 Pt, 16:9）。

关键约束：
- 画布: 960×540 Pt
- 顶栏: 0-44 Pt（深蓝底 + 章节名）
- 内容区: Y=58..486, X=40..920
- 底栏: 494-540 Pt（深蓝底 + 白字金句 + 红色关键词）
"""

from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

from . import theme

from pptx.oxml.ns import qn
from lxml import etree

from pathlib import Path
PPT_ROOT = Path(__file__).resolve().parent.parent


def set_run_font(run, role: str = "body"):
    """统一微软雅黑。"""
    run.font.name = theme.FONT_BODY
    rPr = run._r.get_or_add_rPr()
    for ea in rPr.findall(qn('a:ea')):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', theme.FONT_BODY)


def add_crest(slide):
    """顶部条右侧校徽。"""
    try:
        crest_path = PPT_ROOT / "assets" / "图片1.jpg"
        if not crest_path.exists():
            return
        slide.shapes.add_picture(
            str(crest_path),
            left=theme.SLIDE_WIDTH - theme.CREST_MARGIN_R - theme.CREST_WIDTH,
            top=theme.CREST_MARGIN_T,
            width=theme.CREST_WIDTH,
            height=theme.CREST_HEIGHT,
        )
    except Exception:
        pass


def apply_chrome_v2(slide, chapter_idx: int, page_num: int):
    """顶部深蓝条 + 底部页码条。"""
    chapter = theme.CHAPTERS[chapter_idx - 1]

    # 顶部深蓝条
    add_rect(slide, left=0, top=0,
             width=theme.SLIDE_WIDTH, height=theme.HEADER_BAR_HEIGHT,
             fill=theme.NAVY)

    # 左上章节信息
    tb = slide.shapes.add_textbox(left=Pt(12), top=Pt(10),
                                   width=Pt(700), height=Pt(24))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = chapter["num"] + "    "
    run1.font.size = Pt(theme.HEADER_CHAPTER_FONT_SIZE)
    run1.font.bold = True
    run1.font.color.rgb = hex_to_rgb(theme.WHITE)
    set_run_font(run1)
    run2 = p.add_run()
    run2.text = chapter["title"]
    run2.font.size = Pt(theme.HEADER_CHAPTER_FONT_SIZE)
    run2.font.color.rgb = hex_to_rgb("#CCCCCC")
    set_run_font(run2)

    # 校徽
    add_crest(slide)

    # 底部深蓝条（仅页码）
    add_rect(slide, left=0, top=theme.SLIDE_HEIGHT - Pt(16),
             width=theme.SLIDE_WIDTH, height=Pt(16),
             fill=theme.NAVY)
    add_textbox(slide,
                left=theme.SLIDE_WIDTH - Pt(80),
                top=theme.SLIDE_HEIGHT - Pt(14),
                width=Pt(66), height=Pt(12),
                text=f"{page_num} / {theme.TOTAL_PAGES}",
                font_size=8, color="#AAAAAA",
                align=PP_ALIGN.RIGHT)


def add_bottom_bar(slide, text: str, *, highlight_words: list = None):
    """
    底部"金句"条 — 参考路演PPT核心特征。
    深蓝底 + 白色大字 + 红色仅高亮关键词。
    位置固定在 Y=494..540（46 Pt 高）。
    """
    bar_top = Pt(494)
    bar_h = Pt(46)

    add_rect(slide, left=0, top=bar_top,
             width=theme.SLIDE_WIDTH, height=bar_h,
             fill=theme.NAVY)

    tb = slide.shapes.add_textbox(
        left=theme.MARGIN_LR, top=bar_top + Pt(6),
        width=theme.SLIDE_WIDTH - 2 * theme.MARGIN_LR, height=bar_h - Pt(12),
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    if highlight_words:
        import re
        pattern = '|'.join(re.escape(w) for w in highlight_words)
        parts = re.split(f'({pattern})', text)
        for part in parts:
            run = p.add_run()
            run.text = part
            run.font.size = Pt(theme.FONT_SIZES["bottom_bar"])
            run.font.bold = True
            if part in (highlight_words or []):
                run.font.color.rgb = hex_to_rgb(theme.RED)
            else:
                run.font.color.rgb = hex_to_rgb(theme.WHITE)
            set_run_font(run)
    else:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(theme.FONT_SIZES["bottom_bar"])
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(theme.WHITE)
        set_run_font(run)

    return tb


def hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def add_textbox(slide, left, top, width, height, text, *,
                font_size=14, bold=False, color=theme.DARK_TEXT,
                font_name=theme.FONT_BODY, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, fill=None):
    """通用文本框。"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = hex_to_rgb(fill)
        tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color)
    set_run_font(run)
    return tb


def add_rect(slide, left, top, width, height, *,
             fill=theme.NAVY, line=None, line_width=0):
    """通用矩形。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill)
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = hex_to_rgb(line)
        shp.line.width = Pt(line_width)
    shp.shadow.inherit = False
    return shp


def add_page_title(slide, title: str, subtitle: str = None,
                   top: int = Pt(54), *,
                   accent_color=theme.NAVY):
    """
    内容页标题：蓝色短横线 + 深蓝大标题 + 灰色副标题。
    返回正文起始 top（≈122pt）。
    """
    # 蓝色短横线
    add_rect(slide, left=theme.MARGIN_LR, top=top,
             width=Pt(36), height=Pt(3), fill=accent_color)
    # 大标题
    add_textbox(slide, left=theme.MARGIN_LR, top=top + Pt(10),
                width=theme.SLIDE_WIDTH - 2 * theme.MARGIN_LR, height=Pt(42),
                text=title, font_size=theme.FONT_SIZES["page_title"],
                bold=True, color=theme.NAVY)
    if subtitle:
        add_textbox(slide, left=theme.MARGIN_LR, top=top + Pt(54),
                    width=theme.SLIDE_WIDTH - 2 * theme.MARGIN_LR, height=Pt(24),
                    text=subtitle, font_size=theme.FONT_SIZES["subtitle"],
                    color=theme.TEXT_MUTED)
        return top + Pt(86)
    return top + Pt(64)


def build_section_divider(prs, chapter_idx: int) -> int:
    """
    章节分隔页：深蓝全幅背景 + 巨大半透明数字 + 白色标题 + 要点预览。
    """
    chapter = theme.CHAPTERS[chapter_idx - 1]
    page_num = len(prs.slides) + 1

    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 全幅深蓝背景
    add_rect(slide, left=Pt(0), top=Pt(0),
             width=theme.SLIDE_WIDTH, height=theme.SLIDE_HEIGHT,
             fill=theme.NAVY)

    # 顶部条
    apply_chrome_v2(slide, chapter_idx, page_num)

    # 巨大数字（半透明蓝色）
    add_textbox(slide, left=Pt(48), top=Pt(70),
                width=Pt(420), height=Pt(380),
                text=chapter["num"], font_size=200, bold=True,
                color="#304FA0")

    # 章节中文名
    add_textbox(slide, left=Pt(480), top=Pt(160),
                width=Pt(440), height=Pt(60),
                text=chapter["title"], font_size=34, bold=True,
                color=theme.WHITE)

    # 红色短横线（唯一红色点缀）
    add_rect(slide, left=Pt(480), top=Pt(230),
             width=Pt(40), height=Pt(3),
             fill=theme.RED)

    # 英文
    en_titles = {
        "01": "Introduction & Requirements",
        "02": "System Architecture & Technology",
        "03": "Five Intelligent Agents",
        "04": "Key Technologies & Outlook",
        "05": "Acknowledgement",
    }
    add_textbox(slide, left=Pt(480), top=Pt(250),
                width=Pt(440), height=Pt(26),
                text=en_titles[chapter["num"]],
                font_size=15, color="#8899BB")

    # 要点预览
    chapter_highlights = {
        "01": ["AI教育市场趋势与差距分析", "传统平台3大痛点", "本项目4大差异化定位", "A3赛题全覆盖对标"],
        "02": ["四层轻量化架构设计", "前端+AI+可视化技术选型", "零外部依赖一行启动"],
        "03": ["画像·资源·路径·辅导·评估", "5类智能体职责与协作", "每个智能体的核心能力拆解"],
        "04": ["多智能体协同框架", "流式交互与思考可视化", "路径与练习双向同步", "5项创新+3大未来方向"],
        "05": ["项目成果总结", "开源组件致谢", "Q&A交流"],
    }
    highlights = chapter_highlights.get(chapter["num"], [])
    hl_top = Pt(300)
    for i, line in enumerate(highlights):
        add_textbox(slide, left=Pt(480), top=hl_top + i * Pt(32),
                    width=Pt(440), height=Pt(26),
                    text=f"·  {line}", font_size=14, color="#CCDDFF")

    add_textbox(slide, left=Pt(480), top=hl_top + len(highlights) * Pt(32) + Pt(16),
                width=Pt(440), height=Pt(20),
                text=f"本章范围 · P. {chapter['pages']}",
                font_size=12, color="#8899BB")

    return page_num
