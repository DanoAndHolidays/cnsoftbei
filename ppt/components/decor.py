"""
PPT 装饰辅助：渐变背景、装饰几何、Callout、KPI Badge。

所有新辅助都集中在这里，方便改主题。
"""

from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

from . import theme
from .layout import hex_to_rgb, add_textbox, add_rect
from .shapes import add_card, add_color_block


# ============ 渐变背景（XML 注入）============
# python-pptx 不直接支持渐变填充，需手动注入 <a:gradFill>。

def add_gradient_rect(slide, left, top, width, height, *,
                      color1=theme.PRIMARY, color2=theme.PRIMARY_DARK,
                      direction="horizontal"):
    """
    双色渐变矩形。

    direction:
      - "horizontal": 左 color1 → 右 color2
      - "vertical":   上 color1 → 下 color2
      - "diagonal":   左上 color1 → 右下 color2
    """
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shp.line.fill.background()
    shp.shadow.inherit = False

    # 1. 清空 fill（移除默认 noFill / solidFill）
    sp = shp.fill._xPr
    # 移除所有 <a:*Fill> 子元素
    for tag in ("a:noFill", "a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill"):
        for el in sp.findall(qn(tag)):
            sp.remove(el)

    # 2. 构造 gradFill
    grad = etree.SubElement(sp, qn("a:gradFill"), rotWithShape="1")

    # 方向 → 角度
    angle = {
        "horizontal": 0,
        "vertical": 5400000,   # 90°
        "diagonal": 13500000,  # 135°
    }.get(direction, 0)

    gsLst = etree.SubElement(grad, qn("a:gsLst"))
    pos1, pos2 = ("0", "100000")
    for pos, color in [(pos1, color1), (pos2, color2)]:
        gs = etree.SubElement(gsLst, qn("a:gs"), pos=pos)
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", color.lstrip("#").upper())

    etree.SubElement(grad, qn("a:lin"), ang=str(angle), scaled="1")
    etree.SubElement(grad, qn("a:tileRect"))

    return shp


# ============ 装饰几何 ============

def add_corner_accent(slide, corner="top-right", *, color=theme.PRIMARY, size=Pt(40)):
    """
    在指定角加一个 L 形装饰（两条 4pt 粗线）。

    corner: top-left | top-right | bottom-left | bottom-right
    """
    slide_w = theme.SLIDE_WIDTH
    slide_h = theme.SLIDE_HEIGHT
    if corner == "top-right":
        h_x, h_y = slide_w - size, Pt(50)
        v_x, v_y = slide_w - Pt(4), Pt(50)
        h_w, h_h = size, Pt(4)
        v_w, v_h = Pt(4), size
    elif corner == "top-left":
        h_x, h_y = Pt(0), Pt(50)
        v_x, v_y = Pt(0), Pt(50)
        h_w, h_h = size, Pt(4)
        v_w, v_h = Pt(4), size
    elif corner == "bottom-right":
        h_x, h_y = slide_w - size, slide_h - Pt(54)
        v_x, v_y = slide_w - Pt(4), slide_h - Pt(54) - size
        h_w, h_h = size, Pt(4)
        v_w, v_h = Pt(4), size
    elif corner == "bottom-left":
        h_x, h_y = Pt(0), slide_h - Pt(54)
        v_x, v_y = Pt(0), slide_h - Pt(54) - size
        h_w, h_h = size, Pt(4)
        v_w, v_h = Pt(4), size
    else:
        raise ValueError(f"unknown corner: {corner}")

    add_rect(slide, h_x, h_y, h_w, h_h, fill=color)
    add_rect(slide, v_x, v_y, v_w, v_h, fill=color)


def add_dot_grid(slide, left, top, width, height, *,
                 color=theme.PRIMARY_LIGHT, spacing=Pt(20), dot_size=Pt(3)):
    """
    在指定区域平铺浅色圆点（适合做背景装饰）。
    """
    n_cols = int(width / spacing)
    n_rows = int(height / spacing)
    for r in range(n_rows):
        for c in range(n_cols):
            cx = left + c * spacing + dot_size // 2
            cy = top + r * spacing + dot_size // 2
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, dot_size, dot_size)
            dot.fill.solid()
            dot.fill.fore_color.rgb = hex_to_rgb(color)
            dot.line.fill.background()
            dot.shadow.inherit = False


# ============ 重点高亮 Callout ============

def add_callout_box(slide, left, top, width, height, *,
                    icon="💡", title="", body="", color=theme.PRIMARY,
                    title_size=16, body_size=12):
    """
    高亮 Callout：左侧色条 + emoji icon + 加粗标题 + 正文描述。

    用于突出"创新点"、"关键设计"等需要被注意到的内容。
    """
    add_card(slide, left, top, width, height, fill=theme.ACCENT_BG, border=color, border_width=1.5)
    add_color_block(slide, left, top, Pt(6), height, color)

    # emoji icon
    if icon:
        add_textbox(
            slide, left + Pt(16), top + Pt(10), Pt(32), Pt(28),
            text=icon, font_size=18, color=color, bold=True,
        )

    # title
    if title:
        add_textbox(
            slide, left + Pt(52), top + Pt(10), width - Pt(64), Pt(28),
            text=title, font_size=title_size, bold=True, color=theme.PRIMARY_DARK,
        )

    # body
    if body:
        body_top = top + Pt(42) if title else top + Pt(10)
        body_h = height - (body_top - top) - Pt(8)
        add_textbox(
            slide, left + Pt(16), body_top, width - Pt(28), body_h,
            text=body, font_size=body_size, color=theme.TEXT,
        )


def add_kpi_badge(slide, left, top, width, height, *,
                  value="", label="", color=theme.PRIMARY,
                  value_size=36, label_size=12):
    """
    KPI 数字徽章：大数字 + 小标签 + 左侧色条。

    用于总结页数字（5/7/12/19...）。
    """
    add_card(slide, left, top, width, height, fill=theme.ACCENT_BG, border=color, border_width=1.0)
    add_color_block(slide, left, top, Pt(5), height, color)

    add_textbox(
        slide, left + Pt(14), top + Pt(8), width // 2 - Pt(14), height - Pt(16),
        text=value, font_size=value_size, bold=True, color=color,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_textbox(
        slide, left + width // 2, top + Pt(8), width // 2 - Pt(8), height - Pt(16),
        text=label, font_size=label_size, color=theme.TEXT_MUTED,
        anchor=MSO_ANCHOR.MIDDLE,
    )
