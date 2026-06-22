"""
PPT 暗色代码块：4 色简易语法着色（关键字 / 字符串 / 注释 / 文本）。
"""

import re
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from . import theme
from .layout import hex_to_rgb, add_textbox, add_rect


# 颜色定义
COLOR_KEYWORD  = "#C586C0"
COLOR_STRING   = "#CE9178"
COLOR_COMMENT  = "#6A9955"
COLOR_NORMAL   = "#D4D4D4"
COLOR_FUNCTION = "#DCDCAA"
COLOR_NUMBER   = "#B5CEA8"


# 关键字集合
KEYWORDS = {
    "def", "class", "import", "from", "return", "if", "else", "elif",
    "for", "while", "in", "of", "function", "const", "let", "var",
    "new", "this", "async", "await", "try", "catch", "throw", "yield",
    "interface", "type", "extends", "implements", "export", "default",
    "public", "private", "protected", "static", "void", "null", "true",
    "false", "undefined", "None", "True", "False",
}


def tokenize(line: str):
    """
    把一行代码切成 (text, type) 列表。
    """
    tokens = []
    i = 0
    n = len(line)

    while i < n:
        ch = line[i]

        # 注释
        if ch == "#" and (i == 0 or line[i-1] != ":"):
            tokens.append((line[i:], "COMMENT"))
            return tokens
        if ch == "/" and i + 1 < n and line[i+1] == "/":
            tokens.append((line[i:], "COMMENT"))
            return tokens

        # 字符串
        if ch in ('"', "'", "`"):
            quote = ch
            j = i + 1
            while j < n and line[j] != quote:
                if line[j] == "\\" and j + 1 < n:
                    j += 2
                else:
                    j += 1
            j = min(j + 1, n)
            tokens.append((line[i:j], "STRING"))
            i = j
            continue

        # 数字
        if ch.isdigit():
            j = i
            while j < n and (line[j].isdigit() or line[j] == "."):
                j += 1
            tokens.append((line[i:j], "NUMBER"))
            i = j
            continue

        # 标识符
        if ch.isalpha() or ch == "_":
            j = i
            while j < n and (line[j].isalnum() or line[j] == "_"):
                j += 1
            word = line[i:j]
            if word in KEYWORDS:
                k = j
                while k < n and line[k] == " ":
                    k += 1
                if k < n and line[k] == "(":
                    tokens.append((word, "FUNCTION"))
                else:
                    tokens.append((word, "KEYWORD"))
            else:
                tokens.append((word, "NORMAL"))
            i = j
            continue

        tokens.append((ch, "NORMAL"))
        i += 1

    return tokens


def render_code_block(slide, left, top, width, height, code: str, *,
                      font_size=11, line_spacing=1.25, lang_label=None):
    """在指定区域绘制暗色代码块"""
    bg = add_rect(slide, left, top, width, height, fill="#1E1E1E")
    bg.line.color.rgb = hex_to_rgb("#333333")
    bg.line.width = Pt(0.75)
    bg.shadow.inherit = False

    if lang_label:
        label_w = Pt(60)
        label_h = Pt(20)
        add_rect(slide, left=left, top=top, width=label_w, height=label_h, fill="#3C3C3C")
        add_textbox(
            slide, left=left, top=top, width=label_w, height=label_h,
            text=lang_label, font_size=9, color="#CCCCCC",
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

    tb = slide.shapes.add_textbox(
        left + Pt(12), top + (Pt(26) if lang_label else Pt(8)),
        width - Pt(24), height - (Pt(34) if lang_label else Pt(16)),
    )
    tb.fill.background()
    tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)

    color_map = {
        "KEYWORD":  COLOR_KEYWORD,
        "STRING":   COLOR_STRING,
        "COMMENT":  COLOR_COMMENT,
        "NORMAL":   COLOR_NORMAL,
        "FUNCTION": COLOR_FUNCTION,
        "NUMBER":   COLOR_NUMBER,
    }

    lines = code.split("\n")
    for idx, line in enumerate(lines):
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        if not line:
            r = p.add_run()
            r.text = " "
            r.font.name = theme.FONT_MONO
            r.font.size = Pt(font_size)
            continue
        for text, ttype in tokenize(line):
            r = p.add_run()
            r.text = text
            r.font.name = theme.FONT_MONO
            r.font.size = Pt(font_size)
            r.font.italic = (ttype == "COMMENT")
            r.font.color.rgb = hex_to_rgb(color_map[ttype])