"""
PPT 流程图：节点（圆角矩形）+ 箭头（连接线）。

提供 build_node 和 build_arrow 两个原语；具体流程图布局由 slide 模块自行组织。
"""

from pptx.util import Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from . import theme
from .layout import hex_to_rgb, add_textbox


def build_node(slide, left, top, width, height, text: str, *,
               fill=theme.PRIMARY, text_color=theme.WHITE,
               font_size=12, bold=True, radius=0.15):
    """流程图节点：圆角矩形 + 居中文字"""
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb(fill)
    shp.line.fill.background()
    shp.adjustments[0] = radius
    shp.shadow.inherit = False

    tf = shp.text_frame
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = theme.FONT_FAMILY
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(text_color)
    return shp


def build_arrow(slide, x1, y1, x2, y2, *,
                color=theme.TEXT_MUTED, width_pt=1.5, dashed=False):
    """
    直线箭头连接器。从 (x1,y1) 指向 (x2,y2)。
    x1/y1/x2/y2 都是 EMU 整数。
    """
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = hex_to_rgb(color)
    line.line.width = Pt(width_pt)
    if dashed:
        # python-pptx 设置 dash style: 7 = dash
        try:
            line.line._get_or_add_ln()
            from pptx.oxml.ns import qn
            ln = line.line._get_or_add_ln()
            prstDash = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
            ln.append(prstDash)
        except Exception:
            pass
    # 末端箭头
    line.line._get_or_add_ln()
    from pptx.oxml.ns import qn
    ln = line.line._get_or_add_ln()
    tail = ln.find(qn("a:tailEnd"))
    if tail is None:
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
    return line


def build_horizontal_pipeline(slide, top, left_start, total_width, height, items: list[str], *,
                              fill=theme.PRIMARY, font_size=12, gap_pt=12):
    """
    一行水平流程：items 中的每个文字顺次排开，相邻之间用箭头连接。

    返回每个节点的 (left, top, width, height)，方便 slide 做后续标注。
    """
    n = len(items)
    gap = Pt(gap_pt)
    node_w = (total_width - gap * (n - 1)) // n
    positions = []
    for i, text in enumerate(items):
        x = left_start + i * (node_w + gap)
        shp = build_node(slide, x, top, node_w, height, text, fill=fill, font_size=font_size)
        positions.append((x, top, node_w, height))
        # 在节点右侧画箭头（除最后一个外）
        if i < n - 1:
            arrow_x1 = x + node_w
            arrow_x2 = x + node_w + gap
            arrow_y = top + height // 2
            build_arrow(slide, arrow_x1, arrow_y, arrow_x2, arrow_y)
    return positions