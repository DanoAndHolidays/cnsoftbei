"""
第 15 页 · 关键技术 3：路径 ↔ 练习双向同步。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from components.code_block import render_code_block
from components.flow_diagram import build_node, build_arrow


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=4, page_num=15)

    add_color_block(slide, Pt(80), Pt(70), Pt(140), Pt(28), "#722ED1")
    add_textbox(slide, Pt(80), Pt(70), Pt(140), Pt(28),
                text="关键技术 03", font_size=14, bold=True, color=theme.WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left=Pt(232), top=Pt(70), width=Pt(900), height=Pt(38),
                text="路径 ↔ 练习 双向同步",
                font_size=24, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, left=Pt(232), top=Pt(112), width=Pt(900), height=Pt(20),
                text="StructuredLearningNode + localStorage 事件 + 80% 完成度阈值",
                font_size=12, color=theme.TEXT_MUTED)

    # 上半：数据流图（左：Path → 中：localStorage → 右：Practice → 回到 Path）
    add_textbox(slide, Pt(80), Pt(170), Pt(1080), Pt(24),
                text="▶ 数据流图", font_size=14, bold=True, color="#722ED1")

    nodes_y = Pt(220)
    node_h = Pt(70)
    node_w = Pt(180)
    # 4 个节点：Path 页 / localStorage / Practice 页 / Assessment 页
    positions = [
        (Pt(80),  "Path 页",       theme.PRIMARY,  "写入 activeStructuredPath"),
        (Pt(360), "localStorage",  "#FA8C16",     "单一数据源"),
        (Pt(640), "Practice 页",   theme.SUCCESS,  "读取 + 过滤模块"),
        (Pt(920), "Assessment 页", "#722ED1",     "读取 + 进度展示"),
    ]
    for x, text, color, _ in positions:
        build_node(slide, x, nodes_y, node_w, node_h, text, fill=color, font_size=14)

    # 节点间箭头
    for i in range(3):
        x1 = positions[i][0] + node_w
        x2 = positions[i+1][0]
        y = nodes_y + node_h // 2
        build_arrow(slide, x1, y, x2, y, color="#722ED1", width_pt=2.0)

    # 节点说明
    for x, _, _, desc in positions:
        add_textbox(slide, x, nodes_y + node_h + Pt(8), node_w, Pt(20),
                    text=desc, font_size=10, color=theme.TEXT_MUTED, align=PP_ALIGN.CENTER)

    # 循环箭头（从 Assessment 回到 Path）
    loop_x = Pt(1020)
    loop_y_start = nodes_y + node_h + Pt(8)
    loop_y_end = nodes_y - Pt(10)
    build_arrow(slide, loop_x, loop_y_start, loop_x, loop_y_end, color="#13C2C2", width_pt=1.5, dashed=True)
    add_textbox(slide, Pt(800), nodes_y + node_h + Pt(30), Pt(280), Pt(20),
                text="做题结果回流 → 画像更新 → 推荐新路径",
                font_size=10, color="#13C2C2", align=PP_ALIGN.CENTER)

    # 下半：核心数据结构 + 关键阈值
    code = """interface StructuredLearningNode {
  id: string;
  title: string;
  questionBankId: string;   // ← 关键：关联题库
  moduleId: string;         // ← 关键：关联模块
  moduleName?: string;
  isEntry?: boolean;
  valid?: boolean;          // ← 引用校验结果
}

const THRESHOLD = 0.8;  // ← 完成度阈值

// 做完题后 dispatch 事件
window.dispatchEvent(
  new CustomEvent('moduleProgressUpdated', {
    detail: { moduleId, correctRate }
  })
);"""
    render_code_block(slide, Pt(80), Pt(390), Pt(600), Pt(220), code, font_size=10, lang_label="TS")

    # 右下：3 个关键事实
    facts_x = Pt(720)
    facts_w = Pt(490)
    facts = [
        ("12 题库 × 48 题", "576 道题全预写，AI 只判分不生成", theme.PRIMARY),
        ("80% 阈值", "模块完成度达 80% 自动标记为 done", "#FA8C16"),
        ("customEvent", "解耦通信，跨页面 / 跨组件零依赖", "#722ED1"),
    ]
    fact_h = Pt(70)
    for i, (title, desc, color) in enumerate(facts):
        y = Pt(390) + i * (fact_h + Pt(10))
        add_card(slide, facts_x, y, facts_w, fact_h, fill=theme.ACCENT_BG, border=color, border_width=1.0)
        add_textbox(slide, facts_x + Pt(16), y + Pt(8), Pt(140), Pt(28),
                    text=title, font_size=14, bold=True, color=color)
        add_textbox(slide, facts_x + Pt(160), y + Pt(12), facts_w - Pt(180), Pt(50),
                    text=desc, font_size=11, color=theme.TEXT)