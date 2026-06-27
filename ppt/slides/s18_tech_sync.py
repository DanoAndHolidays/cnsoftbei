"""
第 18 页 · 路径↔练习双向同步 · 学术商务版。
上：数据流  |  下左：接口+DAG  |  下右：3事实卡  +  底部金句。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from components.mini_diagram import render_ascii_block
from components.flow_diagram import build_node, build_arrow


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=4, page_num=18)

    add_page_title(slide, "路径与练习双向同步闭环",
                   subtitle="StructuredLearningNode绑定题库/模块ID · localStorage单一数据源 · 80%阈值标记",
                   accent_color=theme.NAVY)

    LX, LW = Pt(40), Pt(440)
    RX, RW = Pt(510), Pt(410)

    # ── 上半：数据流图 ──
    add_textbox(slide, Pt(40), Pt(130), Pt(840), Pt(22),
                text="数据流闭环：Path → localStorage → Practice → Assessment",
                font_size=15, bold=True, color=theme.NAVY)

    flow_top = Pt(158)
    node_h = Pt(46)
    node_w = Pt(180)
    positions = [
        (Pt(40), "Path 页", theme.NAVY, "激活路径写入"),
        (Pt(252), "localStorage", theme.BLUE_MID, "单一数据源"),
        (Pt(464), "Practice 页", theme.NAVY, "读取+过滤+判分"),
        (Pt(676), "Assessment", theme.BLUE_MID, "读取+进度展示"),
    ]
    y_arrow = flow_top + node_h // 2
    for i, (x, text, color, desc) in enumerate(positions):
        build_node(slide, x, flow_top, node_w, node_h, text, fill=color, font_size=12)
        if i < 3:
            x1 = x + node_w
            x2 = positions[i + 1][0]
            build_arrow(slide, x1, y_arrow, x2, y_arrow, color=theme.NAVY, width_pt=1.5)

    for x, _, color, desc in positions:
        add_textbox(slide, x, flow_top + node_h + Pt(4), node_w, Pt(18),
                    text=desc, font_size=10, color=color, align=PP_ALIGN.CENTER)

    add_textbox(slide, Pt(120), flow_top + node_h + Pt(24), Pt(680), Pt(20),
                text="🔄 做题→画像更新→推荐新路径→新练习→评估刷新（完整学习闭环）",
                font_size=12, bold=True, color=theme.NAVY, align=PP_ALIGN.CENTER)

    # ── 下半左 ──
    flow_end = flow_top + node_h + Pt(50)
    add_textbox(slide, LX, flow_end, LW, Pt(22),
                text="StructuredLearningNode 接口 + 依赖DAG", font_size=14, bold=True, color=theme.NAVY)

    code = """interface StructuredLearningNode {
    id: string;  title: string
    questionBankId: string   // 绑定题库
    moduleId: string         // 绑定模块
    score?: number           // 完成度 0-1
    dependencies?: string[]  // 前置节点
}
const THRESHOLD = 0.8  // 80%自动标记

[基础语法]→[控制流]→[函数]→[装饰器]
   ↓         ↓        ↓
[数据类型][异常处理][闭包/迭代器]"""
    render_ascii_block(slide, LX, flow_end + Pt(30), LW, Pt(144), code,
                       title="TypeScript + DAG", title_color=theme.NAVY,
                       font_size=10, border_color=theme.NAVY, fill=theme.LIGHT_GRAY)

    # ── 下半右：3 事实卡 ──
    facts = [
        ("12题库×48题 = 576题全预写",
         "全部预写，AI只判分不生成题目；从根源杜绝AI幻觉；题型比例6:3:1科学分布",
         theme.NAVY),
        ("80%完成度阈值自动标记",
         "模块完成≥0.8→自动标记done；customEvent广播事件；Path Banner+雷达图实时刷新",
         theme.BLUE_MID),
        ("customEvent跨页面解耦通信",
         "不用Redux/Zustand；仅用浏览器原生CustomEvent；零依赖·跨页面·实时性<100ms",
         theme.NAVY),
    ]
    card_h = Pt(58)
    for i, (title, desc, color) in enumerate(facts):
        y = flow_end + i * (card_h + Pt(6))
        add_card(slide, RX, y, RW, card_h, fill=theme.LIGHT_GRAY, border=color, border_width=1.0)
        add_rect(slide, RX, y, Pt(3), card_h, fill=color)
        add_textbox(slide, RX + Pt(14), y + Pt(4), RW - Pt(28), Pt(22),
                    text=title, font_size=13, bold=True, color=color)
        add_textbox(slide, RX + Pt(14), y + Pt(28), RW - Pt(28), Pt(28),
                    text=desc, font_size=11, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "双向同步是系统的\"神经网络\"：路径定义→练习执行→进度反馈→路径调整，全程自动化",
                   highlight_words=["神经网络", "全程自动化"])
