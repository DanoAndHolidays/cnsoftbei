"""
第 8 页 · 画像构建智能体（Profile Agent）🟠
"""

import os
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from PIL import Image


def _screenshot(slide, img_rel_path: str, left, top, width, height):
    """等比缩放插入截图（来自 assets/screenshot/）"""
    ppt_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    img_path = os.path.abspath(os.path.join(ppt_dir, theme.SCREENSHOT_DIR, img_rel_path))
    if not os.path.exists(img_path):
        # fallback: 占位卡
        add_card(slide, left, top, width, height, fill=theme.ACCENT_BG, border=theme.BORDER)
        add_textbox(slide, left, top, width, height,
                    text=f"[截图缺失: {img_rel_path}]", font_size=12, color=theme.TEXT_SUBTLE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    # 等比缩放
    img = Image.open(img_path)
    iw, ih = img.size
    ratio = min(width / iw, height / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    x = left + (width - w) // 2
    y = top + (height - h) // 2
    slide.shapes.add_picture(img_path, x, y, width=w, height=h)


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=3, page_num=8)

    color = theme.AGENT_COLORS["profile"]
    name = theme.AGENT_NAMES_CN["profile"]
    en = theme.AGENT_NAMES_EN["profile"]
    emoji = theme.AGENT_EMOJI["profile"]

    # 标题
    add_color_block(slide, Pt(80), Pt(70), Pt(10), Pt(28), color)
    add_textbox(slide, left=Pt(100), top=Pt(70), width=Pt(700), height=Pt(38),
                text=f"{name} {emoji}", font_size=24, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, left=Pt(100), top=Pt(110), width=Pt(600), height=Pt(20),
                text=en, font_size=12, color=theme.TEXT_SUBTLE)

    # 左侧截图
    _screenshot(slide, theme.SCREENSHOTS["profile"], Pt(80), Pt(160), Pt(560), Pt(440))

    # 右侧三段式
    right_x = Pt(680)
    right_w = Pt(540)
    add_card(slide, right_x, Pt(160), right_w, Pt(110), fill=theme.ACCENT_BG, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(170), right_w - Pt(32), Pt(24),
                text="▶ 角色定位", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(196), right_w - Pt(32), Pt(70),
                text="通过对话式交互理解学习者，构建 6 维动态画像；随学习过程持续更新。",
                font_size=12, color=theme.TEXT)

    add_card(slide, right_x, Pt(280), right_w, Pt(160), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(290), right_w - Pt(32), Pt(24),
                text="▶ 核心能力（6 维画像）", font_size=13, bold=True, color=color)
    dims = ["知识基础", "认知风格", "易错偏好", "学习节奏", "兴趣方向", "学习习惯"]
    for i, d in enumerate(dims):
        col = i % 2
        row = i // 2
        x = right_x + Pt(16) + col * Pt(260)
        y = Pt(320) + row * Pt(32)
        add_rect(slide, x, y + Pt(4), Pt(8), Pt(8), fill=color)
        add_textbox(slide, x + Pt(16), y, Pt(230), Pt(20),
                    text=d, font_size=12, color=theme.TEXT)

    add_card(slide, right_x, Pt(450), right_w, Pt(150), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(460), right_w - Pt(32), Pt(24),
                text="▶ 关键特性：随学随新", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(488), right_w - Pt(32), Pt(110),
                text="• 对话式构建：自然语言输入\n• 持久化到 localStorage，跨页面共享\n"
                     "• 做题反馈回流画像（practiceGrader 派发事件）\n• 系统 prompt 注入到所有下游智能体",
                font_size=12, color=theme.TEXT)
