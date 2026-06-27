"""
第 12 页 · 路径规划智能体 · 学术商务版。
"""

from pptx.util import Pt

from components import theme
from components.layout import add_textbox, add_rect, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from slides.s10_agent_profile import _screenshot


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=3, page_num=12)

    color = theme.NAVY
    name = theme.AGENT_NAMES_CN["path"]
    en = theme.AGENT_NAMES_EN["path"]

    add_rect(slide, Pt(40), Pt(56), Pt(32), Pt(3), fill=color)
    add_textbox(slide, left=Pt(40), top=Pt(64), width=Pt(860), height=Pt(38),
                text=f"{name}  ({en})", font_size=26, bold=True, color=theme.DARK_TEXT)
    add_textbox(slide, left=Pt(40), top=Pt(104), width=Pt(860), height=Pt(20),
                text="系统的\"导航层\" — AI自由生成 + 12预定义路径 · 节点绑定题库/模块ID · 80%阈值双向同步",
                font_size=14, color=theme.TEXT_MUTED)

    _screenshot(slide, theme.SCREENSHOTS["path1"], Pt(40), Pt(130), Pt(520), Pt(350))

    rx = Pt(590)
    rw = Pt(330)

    # 卡片 1：角色
    add_card(slide, rx, Pt(130), rw, Pt(65), fill=theme.LIGHT_GRAY, border=color, border_width=0.75)
    add_textbox(slide, rx + Pt(14), Pt(136), rw - Pt(28), Pt(20),
                text="角色定位 — 系统的\"导航层\"", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(158), rw - Pt(28), Pt(30),
                text="根据画像生成结构化路径，将知识点对接到具体题库。支持AI自由生成和12条预定义路径。",
                font_size=12, color=theme.DARK_TEXT)

    # 卡片 2：双轨模式
    add_card(slide, rx, Pt(204), rw, Pt(170), fill=theme.LIGHT_GRAY, border=color, border_width=0.75)
    add_textbox(slide, rx + Pt(14), Pt(210), rw - Pt(28), Pt(20),
                text="双轨模式 + 12条预定义路径", font_size=14, bold=True, color=color)
    tracks = ("【AI 自由生成】输入目标→AI流式输出节点列表\n\n"
              "【12条预定义路径（一键采用）】\n"
              "Python全栈·算法工程师·数据分析师\n"
              "前端开发·后端开发·DevOps工程师\n"
              "机器学习·网络安全·数据库管理员\n"
              "软件测试·系统架构师·全栈Web开发")
    add_textbox(slide, rx + Pt(14), Pt(234), rw - Pt(28), Pt(132),
                text=tracks, font_size=11, color=theme.DARK_TEXT)

    # 卡片 3：双向同步
    add_card(slide, rx, Pt(384), rw, Pt(98), fill=theme.LIGHT_GRAY, border=color, border_width=0.75)
    add_textbox(slide, rx + Pt(14), Pt(390), rw - Pt(28), Pt(20),
                text="路径 ↔ 练习 双向同步闭环", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(412), rw - Pt(28), Pt(64),
                text="· Path激活→localStorage→Practice过滤\n"
                     "· 模块完成≥80%→自动标记done\n"
                     "· customEvent广播→Assessment刷新\n"
                     "· 做题回流画像→推荐新路径→闭环",
                font_size=11, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "路径智能体实现\"学-练-评-荐\"四位一体：从知识点到题库到评估的完整自动化闭环",
                   highlight_words=["学-练-评-荐", "完整自动化"])
