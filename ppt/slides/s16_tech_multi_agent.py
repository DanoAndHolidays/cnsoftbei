"""
第 16 页 · 多智能体协同框架 · 学术商务版。
双列布局 + 底部金句。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from components.mini_diagram import render_ascii_block


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=4, page_num=16)

    add_page_title(slide, "多智能体协同框架与事件总线",
                   subtitle="MultiAgentScheduler + EventEmitter 完全解耦 · 代码精简 40%",
                   accent_color=theme.NAVY)

    LX, LW = Pt(40), Pt(440)
    RX, RW = Pt(510), Pt(410)

    # ── 左上：核心类 ──
    add_textbox(slide, LX, Pt(130), LW, Pt(22),
                text="MultiAgentScheduler 核心类", font_size=15, bold=True, color=theme.NAVY)
    code = """class MultiAgentScheduler {
    agents: Map<Role, Agent>      // 角色映射
    eventBus: EventEmitter        // 事件总线

    registerAgent(role, agent)
    dispatch(role, task) → Result
    runPipeline(tasks) → Result[]
    getAgentStates() → Map
}"""
    render_ascii_block(slide, LX, Pt(156), LW, Pt(140), code,
                       title="TypeScript 核心接口", title_color=theme.NAVY,
                       font_size=11, border_color=theme.NAVY, fill=theme.LIGHT_GRAY)

    # ── 左下：状态机 ──
    add_textbox(slide, LX, Pt(306), LW, Pt(22),
                text="状态机 & 关键设计决策", font_size=15, bold=True, color=theme.NAVY)
    add_card(slide, LX, Pt(332), LW, Pt(158), fill=theme.LIGHT_GRAY, border=theme.NAVY, border_width=1.0)
    add_rect(slide, LX, Pt(332), Pt(3), Pt(158), fill=theme.NAVY)
    add_textbox(slide, LX + Pt(12), Pt(340), LW - Pt(20), Pt(144),
                text="【Planner→Worker 调度模型】\n"
                     "· Planner解析目标→拆解子任务队列\n"
                     "· 6 Worker并行，状态：pending→running→done\n"
                     "· failed自动重试3次（指数退避1s→2s→4s）\n\n"
                     "【事件总线解耦 — 核心决策】\n"
                     "· 5智能体零硬编码依赖，仅通过EventEmitter通信\n"
                     "· 新增智能体只需registerAgent()即可接入\n"
                     "· ResourceGenerator复用调度逻辑\n\n"
                     "效果：代码量-40%，新增智能体成本-80%",
                font_size=12, color=theme.DARK_TEXT)

    # ── 右上：事件总线图 ──
    add_textbox(slide, RX, Pt(130), RW, Pt(22),
                text="5 智能体事件总线协作", font_size=15, bold=True, color=theme.NAVY)
    diagram = """Profile ── Resource ── Path
    │ profile:updated │ resource:done │ path:scheduled
    ▼         ▼         ▼
    ┌─────── EventEmitter (事件总线) ───────┐
    │ 解耦通信：只与事件总线交互，不直接依赖 │
    └───────────────────────────────────────┘
    │         │         │
    ▼         ▼         ▼
Assessment ◄─ Tutor ◄───┘
    │ assessment:done  │ tutor:answered

闭环：做题→profile:updated→推荐新路径→循环"""
    render_ascii_block(slide, RX, Pt(156), RW, Pt(198), diagram,
                       title="事件总线解耦架构", font_size=10,
                       border_color=theme.NAVY, fill=theme.WHITE)

    # ── 右下：协作矩阵 ──
    add_textbox(slide, RX, Pt(366), RW, Pt(22),
                text="5 智能体协作矩阵", font_size=15, bold=True, color=theme.NAVY)
    table = ("智能体      输入            输出         事件\n"
             "────────────────────────────────────────\n"
             "Profile     对话/做题结果   6维画像      profile:updated\n"
             "Resource    学习目标+画像   6类资源      resource:generated\n"
             "Path        画像+目标方向   结构化路径   path:scheduled\n"
             "Tutor       问题+历史+画像  答案+思考    tutor:answered\n"
             "Assessment  practiceState  评估+建议    assessment:done")
    render_ascii_block(slide, RX, Pt(392), RW, Pt(98), table,
                       title="", font_size=10, border_color=theme.NAVY,
                       fill=theme.LIGHT_GRAY)

    add_bottom_bar(slide, "核心创新：5个智能体通过事件总线完全解耦 — 新增智能体只需1行注册代码",
                   highlight_words=["完全解耦", "1行注册代码"])
