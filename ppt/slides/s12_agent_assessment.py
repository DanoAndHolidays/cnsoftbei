"""
第 12 页 · 效果评估智能体（Assessment Agent）💠

由于 Assessment 截图缺失，本页用 components/assessment_mock.py 自绘一张
评估页 mockup 替代。
"""

import os
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from components.assessment_mock import render_assessment_mock


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=3, page_num=12)

    color = theme.AGENT_COLORS["assessment"]
    name = theme.AGENT_NAMES_CN["assessment"]
    en = theme.AGENT_NAMES_EN["assessment"]
    emoji = theme.AGENT_EMOJI["assessment"]

    add_color_block(slide, Pt(80), Pt(70), Pt(10), Pt(28), color)
    add_textbox(slide, left=Pt(100), top=Pt(70), width=Pt(700), height=Pt(38),
                text=f"{name} {emoji}", font_size=24, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, left=Pt(100), top=Pt(110), width=Pt(600), height=Pt(20),
                text=en + "  ·  注：Assessment 截图缺失，此处用 mockup 替代", font_size=12, color=theme.TEXT_SUBTLE)

    # 左侧自绘 mockup
    render_assessment_mock(slide, Pt(80), Pt(160), Pt(620), Pt(500))

    # 右侧说明
    right_x = Pt(740)
    right_w = Pt(480)

    add_card(slide, right_x, Pt(160), right_w, Pt(150), fill=theme.ACCENT_BG, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(170), right_w - Pt(32), Pt(24),
                text="▶ 真实进度同步（不造假）", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(196), right_w - Pt(32), Pt(110),
                text="• 直接读取 practiceState（练习页写入）\n"
                     "• 监听 practiceStateUpdated 事件\n"
                     "• 无练习记录时显示引导提示\n"
                     "• 4 个统计卡 + 模块进度卡全基于真实数据",
                font_size=12, color=theme.TEXT)

    add_card(slide, right_x, Pt(320), right_w, Pt(150), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(330), right_w - Pt(32), Pt(24),
                text="▶ 能力雷达（6 维）", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(358), right_w - Pt(32), Pt(110),
                text="• 知识基础 · 认知风格 · 易错偏好\n"
                     "• 学习节奏 · 兴趣方向 · 学习习惯\n"
                     "• Recharts RadarChart 实时更新",
                font_size=12, color=theme.TEXT)

    add_card(slide, right_x, Pt(480), right_w, Pt(180), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(490), right_w - Pt(32), Pt(24),
                text="▶ 智能调整建议", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(518), right_w - Pt(32), Pt(140),
                text="• 薄弱维度专项练习\n"
                     "• 节奏建议（每日学习时长）\n"
                     "• 兴趣方向 → 相关路径推荐\n"
                     "• 整体调整：放慢/加速/转向",
                font_size=12, color=theme.TEXT)
