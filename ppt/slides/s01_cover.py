"""
第 1 页 · 封面 · 学术商务版。
深蓝全幅 + 白字大标题 + 参赛信息 + 亮点tag。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from components import theme
from components.layout import add_textbox, add_rect, hex_to_rgb


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 深蓝全幅背景
    add_rect(slide, left=Pt(0), top=Pt(0),
             width=theme.SLIDE_WIDTH, height=theme.SLIDE_HEIGHT,
             fill=theme.NAVY)

    # 顶部细红线
    add_rect(slide, left=Pt(0), top=Pt(0),
             width=theme.SLIDE_WIDTH, height=Pt(3),
             fill=theme.RED)

    # 装饰圆 - 左上
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Pt(-60), Pt(40), Pt(260), Pt(260))
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb("#2A4A9A")
    shp.fill.transparency = 0.78
    shp.line.fill.background()

    # 赛题信息
    add_textbox(slide, left=Pt(60), top=Pt(80), width=Pt(800), height=Pt(26),
                text=theme.COVER_INFO["contest"], font_size=16, color="#8899BB")

    # 蓝色短横线
    add_rect(slide, left=Pt(60), top=Pt(120), width=Pt(40), height=Pt(3), fill=theme.RED)

    # 大标题
    add_textbox(slide, left=Pt(60), top=Pt(148), width=Pt(840), height=Pt(110),
                text="学习智能体系统", font_size=60, bold=True, color=theme.WHITE)

    # 副标题
    add_textbox(slide, left=Pt(60), top=Pt(266), width=Pt(840), height=Pt(44),
                text="多智能体协同驱动的个性化学习平台", font_size=24, color=theme.WHITE)

    # 英文
    add_textbox(slide, left=Pt(60), top=Pt(314), width=Pt(840), height=Pt(24),
                text="Multi-Agent Driven Personalized Learning Platform",
                font_size=13, color="#8899BB")

    # 亮点tag行
    highlights = ["5智能体协同", "6维动态画像", "SSE流式交互", "576题库", "92%判分一致率", "零外部依赖"]
    tag_y = Pt(370)
    for i, tag in enumerate(highlights):
        x = Pt(60) + i * Pt(150)
        add_textbox(slide, left=x, top=tag_y, width=Pt(140), height=Pt(22),
                    text=f"✦  {tag}", font_size=13, bold=True, color=theme.RED)

    # 参赛信息
    info_top = Pt(430)
    info_items = [
        ("队伍", theme.COVER_INFO["team_name"]),
        ("学校", theme.COVER_INFO["school"]),
        ("汇报人", theme.COVER_INFO["presenter"]),
        ("指导老师", theme.COVER_INFO["advisor"]),
        ("日期", theme.COVER_INFO["date"]),
    ]
    for i, (label, value) in enumerate(info_items):
        y = info_top + i * Pt(22)
        add_textbox(slide, left=Pt(60), top=y, width=Pt(90), height=Pt(20),
                    text=label, font_size=13, bold=True, color=theme.RED)
        add_textbox(slide, left=Pt(160), top=y, width=Pt(400), height=Pt(20),
                    text=value, font_size=13, color=theme.WHITE)

    # 右下标识
    add_textbox(slide, left=Pt(780), top=Pt(490), width=Pt(160), height=Pt(22),
                text="Learning Agent · A3", font_size=13, bold=True,
                color=theme.RED, align=PP_ALIGN.RIGHT)
