"""
第 5 页 · A3赛题需求对标 · 学术商务版。
左：9项需求  |  右：设计原则 + KPI  +  底部金句。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=1, page_num=5)

    add_page_title(slide, "A3 赛题需求精准对标",
                   subtitle="5 项功能需求 100% 覆盖  ·  4 项非功能需求工程化达标  ·  多项超额完成")

    left_x, right_x = Pt(40), Pt(580)
    left_w, right_w = Pt(520), Pt(340)

    # ── 左侧：功能需求 ──
    add_textbox(slide, left_x, Pt(130), left_w, Pt(24),
                text="5 项功能需求 — 全部超额完成", font_size=16, bold=True, color=theme.NAVY)

    reqs = [
        ("画像构建", "6维对话式动态画像，做题反馈回流，注入下游智能体"),
        ("资源生成", "6类资源多智能体流水线协作，SSE实时状态回传"),
        ("路径规划", "AI自由生成+12预定义路径，节点绑定题库/模块ID"),
        ("智能辅导", "4种模式+追问链+画像注入+缓存去重+流式中断"),
        ("学习评估", "真实数据驱动+6维雷达+进度追踪+智能调整建议"),
    ]
    for i, (title, desc) in enumerate(reqs):
        y = Pt(164) + i * Pt(58)
        add_card(slide, left_x, y, left_w, Pt(49), fill=theme.LIGHT_GRAY)
        add_rect(slide, left_x, y, Pt(4), Pt(49), fill=theme.NAVY)
        add_textbox(slide, left_x + Pt(14), y + Pt(4), Pt(120), Pt(22),
                    text=f"0{i+1}  {title}", font_size=14, bold=True, color=theme.NAVY)
        add_textbox(slide, left_x + Pt(14), y + Pt(26), left_w - Pt(28), Pt(20),
                    text=desc, font_size=12, color=theme.DARK_TEXT)

    # 非功能需求
    nfr_top = Pt(164 + 5 * 58 + 4)
    add_textbox(slide, left_x, nfr_top, left_w, Pt(24),
                text="4 项非功能需求", font_size=15, bold=True, color=theme.RED)

    nfrs = [
        ("流式输出", "SSE打字机效果，首字延迟<2s"),
        ("防幻觉", "576题全预写，AI只判分不生成"),
        ("响应追踪", "3次重试+3min超时+AbortSignal"),
        ("多智能体协作", "事件总线解耦，6Worker并发<8s"),
    ]
    for i, (title, desc) in enumerate(nfrs):
        y = nfr_top + Pt(30) + i * Pt(26)
        add_textbox(slide, left_x + Pt(14), y, Pt(130), Pt(22),
                    text=title, font_size=13, bold=True, color=theme.RED)
        add_textbox(slide, left_x + Pt(150), y, left_w - Pt(164), Pt(22),
                    text=desc, font_size=12, color=theme.DARK_TEXT)

    # ── 右侧：设计原则 ──
    add_textbox(slide, right_x, Pt(130), right_w, Pt(24),
                text="4 条核心设计原则", font_size=16, bold=True, color=theme.NAVY)

    principles = [
        ("个性化优先", "画像驱动所有下游智能体", theme.NAVY),
        ("多智能体协作", "Planner拆任务→Worker并行", theme.NAVY),
        ("流式交互体验", "打字机+<thinking>折叠", theme.NAVY),
        ("数据闭环驱动", "做题→画像→推荐→评估", theme.NAVY),
    ]

    for i, (title, desc, color) in enumerate(principles):
        y = Pt(164) + i * Pt(72)
        add_card(slide, right_x, y, right_w, Pt(62), fill=theme.WHITE, border=color, border_width=1.0)
        add_rect(slide, right_x, y, right_w, Pt(4), fill=color)
        add_textbox(slide, right_x + Pt(12), y + Pt(12), right_w - Pt(24), Pt(24),
                    text=title, font_size=15, bold=True, color=color)
        add_textbox(slide, right_x + Pt(12), y + Pt(36), right_w - Pt(24), Pt(22),
                    text=desc, font_size=12, color=theme.TEXT_MUTED)

    # 核心技术指标
    kpi_top = Pt(164 + 4 * 72 + 10)
    add_textbox(slide, right_x, kpi_top, right_w, Pt(22),
                text="核心技术指标", font_size=15, bold=True, color=theme.NAVY)

    kpis = [("576", "题库总题数"), ("92%", "AI判分一致率"), ("<2s", "流式首字延迟"), ("0", "严重缺陷")]
    kpi_w = Pt(78)
    kpi_h = Pt(52)
    for i, (val, label) in enumerate(kpis):
        x = right_x + i * (kpi_w + Pt(6))
        y = kpi_top + Pt(28)
        add_card(slide, x, y, kpi_w, kpi_h, fill=theme.LIGHT_GRAY, border=theme.NAVY, border_width=0.75)
        add_textbox(slide, x, y + Pt(2), kpi_w, Pt(28),
                    text=val, font_size=22, bold=True, color=theme.NAVY, align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Pt(30), kpi_w, Pt(18),
                    text=label, font_size=10, color=theme.TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ─── 底部金句 ───
    add_bottom_bar(slide, "赛题完成度：功能需求 5/5   非功能需求 4/4   加分项 3/3   综合完成度 100%+",
                   highlight_words=["100%+"])
