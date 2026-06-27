"""
第 17 页 · 流式交互与思考可视化 · 学术商务版。
双列：左函数+时序  |  右3特性卡  +  底部金句。
"""

from pptx.util import Pt

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from components.mini_diagram import render_ascii_block


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=4, page_num=17)

    add_page_title(slide, "流式交互与思考过程可视化",
                   subtitle="SSE + ReadableStream 打字机效果 · 思考过程可折叠 · AbortSignal 中断 · 首字延迟<2s",
                   accent_color=theme.NAVY)

    LX, LW = Pt(40), Pt(440)
    RX, RW = Pt(510), Pt(410)

    # ── 左上：核心函数 ──
    add_textbox(slide, LX, Pt(130), LW, Pt(22),
                text="streamChatCompletion 核心函数", font_size=15, bold=True, color=theme.NAVY)
    code = """async function* streamChatCompletion(
    messages: Message[],
    onChunk: (delta) => void,
    signal?: AbortSignal
): AsyncGenerator {
    const res = await fetch('/anthropic', {
        method:'POST', signal,
        body: JSON.stringify({messages, stream:true})
    })
    const reader = res.body.getReader()
    for await (const chunk of readSSE(reader)) {
        const {text, thinking} = parseChunk(chunk)
        onChunk({text, thinking})
    }
}"""
    render_ascii_block(slide, LX, Pt(156), LW, Pt(170), code,
                       title="TypeScript 核心实现", title_color=theme.NAVY,
                       font_size=11, border_color=theme.NAVY, fill=theme.LIGHT_GRAY)

    # ── 左下：时序 ──
    add_textbox(slide, LX, Pt(336), LW, Pt(22),
                text="SSE 打字机时序 + 取消流程", font_size=15, bold=True, color=theme.NAVY)
    diagram = """用户输入 → POST /anthropic (stream:true)
    │
    ▼
chunk1 → onChunk("你") → UI: "你"
chunk2 → onChunk("好") → UI: "你好"
chunk3 → onChunk("世") → UI: "你好世"
    ... 持续流式直到 stop_reason
    │
    ▼
done → onComplete() → UI: "✓ 完成"

取消：点击取消→abort()→中断→"已取消"
性能：首字<2s · 吞吐>50 tok/s"""
    render_ascii_block(slide, LX, Pt(362), LW, Pt(128), diagram,
                       title="", font_size=10, border_color=theme.NAVY, fill=theme.WHITE)

    # ── 右侧：3 特性卡 ──
    cards = [
        ("打字机效果",
         "fetch+ReadableStream逐chunk解析SSE；\n每收到chunk立即setState更新UI；\n视觉上ChatGPT式逐字出现；\n支持Markdown实时渲染。",
         theme.NAVY),
        ("思考过程可折叠",
         "AI输出<thinking>块自动折叠；\n默认隐藏在\"💭思考过程\"标签下；\n用户点击展开查看AI推理逻辑；\n增强AI可信度与透明度。",
         theme.BLUE_MID),
        ("AbortSignal取消控制",
         "流式中点击红色\"取消\"按钮；\ntrigger abort()→fetch立即中断；\nUI恢复\"已取消\"+取消时间；\n已生成文本保留可继续编辑。",
         theme.NAVY),
    ]
    card_h = Pt(96)
    for i, (title, body, color) in enumerate(cards):
        y = Pt(130) + i * (card_h + Pt(10))
        add_card(slide, RX, y, RW, card_h, fill=theme.WHITE, border=color, border_width=1.5)
        add_rect(slide, RX, y, Pt(4), card_h, fill=color)
        add_textbox(slide, RX + Pt(14), y + Pt(6), RW - Pt(28), Pt(22),
                    text=title, font_size=15, bold=True, color=color)
        add_textbox(slide, RX + Pt(14), y + Pt(32), RW - Pt(28), card_h - Pt(36),
                    text=body, font_size=12, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "流式交互让AI\"开口说话\"：打字机效果 + 思考透明化 = 可信赖的AI学习助手",
                   highlight_words=["开口说话", "透明化", "可信赖"])
