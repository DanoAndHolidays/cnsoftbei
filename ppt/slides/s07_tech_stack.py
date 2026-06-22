"""
第 7 页 · 技术选型。

3 列布局：前端 / AI / 数据可视化。
每列 4-5 个技术栈图标（用纯文字 + 色块代替 logo）+ 选型理由。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=2, page_num=7)

    add_page_title(slide, "技术选型", subtitle="主流成熟框架 + 大模型 API + 轻量数据可视化")

    cols = [
        ("前端框架", theme.PRIMARY, [
            ("React 19",     "最新稳定版，Concurrent Features"),
            ("TypeScript",   "类型安全，IDE 智能提示"),
            ("Vite 5",       "极速冷启动 + HMR"),
            ("Ant Design 6", "企业级 UI 组件库"),
            ("PageCache",    "手动 useState 路由 + useRef 缓存"),
        ]),
        ("AI 层", "#FA8C16", [
            ("大模型 API",     "兼容 Anthropic 协议"),
            ("SSE 流式",       "ReadableStream 打字机效果"),
            ("AbortSignal",    "支持中途取消请求"),
            ("Prompt 工程",    "5 类系统 prompt 注入画像"),
            ("重试 + 超时",     "axios 3 次重试 / 3 分钟超时"),
        ]),
        ("数据可视化", "#13C2C2", [
            ("Recharts 3",   "React 原生 + 声明式 API"),
            ("雷达图",       "6 维度能力评估"),
            ("统计卡片",      "Antd Statistic 组件"),
            ("时间线",        "Antd Timeline 智能建议"),
            ("进度条",        "Antd Progress 模块进度"),
        ]),
    ]
    col_w = Pt(360)
    col_top = Pt(200)
    for ci, (col_name, color, items) in enumerate(cols):
        x = Pt(80) + ci * (col_w + Pt(20))
        # 列标题
        add_card(slide, x, col_top, col_w, Pt(40), fill=color, border=None)
        add_textbox(slide, x, col_top, col_w, Pt(40),
                    text=col_name, font_size=16, bold=True, color=theme.WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # 列表
        for i, (name, desc) in enumerate(items):
            y = col_top + Pt(50) + i * Pt(72)
            add_card(slide, x, y, col_w, Pt(64), fill=theme.ACCENT_BG, border=color, border_width=0.75)
            add_textbox(slide, x + Pt(12), y + Pt(8), col_w - Pt(24), Pt(24),
                        text=name, font_size=14, bold=True, color=theme.PRIMARY_DARK)
            add_textbox(slide, x + Pt(12), y + Pt(34), col_w - Pt(24), Pt(24),
                        text=desc, font_size=11, color=theme.TEXT_MUTED)
