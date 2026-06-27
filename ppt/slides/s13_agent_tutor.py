"""
第 13 页 · 辅导答疑智能体 · 学术商务版。
"""

from pptx.util import Pt

from components import theme
from components.layout import add_textbox, add_rect, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from slides.s10_agent_profile import _screenshot


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=3, page_num=13)

    color = theme.NAVY
    name = theme.AGENT_NAMES_CN["tutor"]
    en = theme.AGENT_NAMES_EN["tutor"]

    add_rect(slide, Pt(40), Pt(56), Pt(32), Pt(3), fill=color)
    add_textbox(slide, left=Pt(40), top=Pt(64), width=Pt(860), height=Pt(38),
                text=f"{name}  ({en})", font_size=26, bold=True, color=theme.DARK_TEXT)
    add_textbox(slide, left=Pt(40), top=Pt(104), width=Pt(860), height=Pt(20),
                text="系统的\"交互层\" — 4 种解答模式 · 追问链深度对话 · 画像注入个性化 · 5 项工程优化",
                font_size=14, color=theme.TEXT_MUTED)

    _screenshot(slide, theme.SCREENSHOTS["tutor"], Pt(40), Pt(130), Pt(520), Pt(350))

    rx = Pt(590)
    rw = Pt(330)

    # 卡片 1：角色 + 4 种模式
    add_card(slide, rx, Pt(130), rw, Pt(148), fill=theme.LIGHT_GRAY)
    add_textbox(slide, rx + Pt(14), Pt(136), rw - Pt(28), Pt(20),
                text="角色定位 + 4 种解答模式", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(156), rw - Pt(28), Pt(18),
                text="一对一智能辅导，根据画像调整回答风格与深度；支持多轮追问链。",
                font_size=12, color=theme.DARK_TEXT)

    modes = [("文字模式", "Markdown+代码高亮"), ("图解模式", "Mermaid/ASCII流程图"),
             ("视频模式", "脚本+时间戳+节奏"), ("代码模式", "可执行+逐行注释")]
    for i, (title, desc) in enumerate(modes):
        col = i % 2
        row = i // 2
        x = rx + Pt(14) + col * Pt(158)
        y = Pt(186) + row * Pt(32)
        add_rect(slide, x, y + Pt(4), Pt(3), Pt(16), fill=color)
        add_textbox(slide, x + Pt(10), y, Pt(88), Pt(20),
                    text=title, font_size=12, bold=True, color=color)
        add_textbox(slide, x + Pt(10), y + Pt(18), Pt(140), Pt(16),
                    text=desc, font_size=10, color=theme.TEXT_MUTED)

    # 卡片 2：5 项工程优化
    add_card(slide, rx, Pt(290), rw, Pt(192), fill=theme.LIGHT_GRAY)
    add_textbox(slide, rx + Pt(14), Pt(296), rw - Pt(28), Pt(20),
                text="5 项深度工程优化", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(318), rw - Pt(28), Pt(158),
                text="① 画像注入System Prompt→自动调整风格难度\n"
                     "② 缓存去重（问题+模式双键）→避免重复调用\n"
                     "③ 追问链parentId/followUpIds→深度多轮对话\n"
                     "④ 点踩重新生成（含原因分析）→全新解答\n"
                     "⑤ AbortSignal取消+红色按钮→随时中断",
                font_size=11, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "辅导智能体将AI辅导从\"冷冰冰的问答\"升级为\"具备记忆与纠错能力的深度互动\"",
                   highlight_words=["记忆", "纠错", "深度互动"])
