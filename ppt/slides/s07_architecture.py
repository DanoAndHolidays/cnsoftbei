"""
第 7 页 · 总体架构 · 学术商务版。
4 层架构纵向展开 + 底部金句。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from components.flow_diagram import build_arrow


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=2, page_num=7)

    add_page_title(slide, "四层轻量化总体架构",
                   subtitle="表现层 / API网关 / 多智能体框架 / 数据层  ·  零外部依赖，一行命令启动")

    layers = [
        ("L1 表现层 (Presentation)", "React 19 + TypeScript + Vite 5 + Ant Design 6",
         "手动路由切换（7页面）· PageCacheContext 状态缓存 · MarkdownRenderer", theme.NAVY),
        ("L2 API 网关 (Gateway)", "Vite Dev Proxy → api.minimaxi.com",
         "Anthropic 协议兼容 · SSE 流式传输 · axios 重试 · AbortSignal 取消", theme.BLUE_MID),
        ("L3 多智能体框架", "MultiAgentScheduler + EventEmitter 事件总线",
         "registerAgent() 注册 · runPipeline() 流水线 · dispatch() 单智能体执行", theme.NAVY),
        ("L4 数据层 (Data)", "localStorage 持久化 + 题库 JSON（12库576题）",
         "studentProfile · practiceState · activeStructuredPath · QAHistory", theme.BLUE_MID),
    ]

    layer_top = Pt(130)
    layer_h = Pt(78)
    layer_gap = Pt(10)

    for i, (name, tech, detail, color) in enumerate(layers):
        y = layer_top + i * (layer_h + layer_gap)
        add_rect(slide, Pt(40), y, Pt(4), layer_h, fill=color)
        add_card(slide, Pt(52), y, Pt(870), layer_h, fill=theme.LIGHT_GRAY, border=None)

        add_textbox(slide, Pt(68), y + Pt(6), Pt(500), Pt(26),
                    text=name, font_size=18, bold=True, color=color)
        add_textbox(slide, Pt(68), y + Pt(30), Pt(600), Pt(22),
                    text=tech, font_size=13, color=theme.DARK_TEXT)
        add_textbox(slide, Pt(68), y + Pt(52), Pt(600), Pt(22),
                    text=detail, font_size=12, color=theme.TEXT_MUTED)

        add_textbox(slide, Pt(830), y + Pt(8), Pt(70), layer_h - Pt(16),
                    text=f"L{i + 1}", font_size=32, bold=True,
                    color=color, align=PP_ALIGN.RIGHT)

        if i < len(layers) - 1:
            arrow_y = y + layer_h + Pt(1)
            build_arrow(slide, Pt(490), arrow_y, Pt(490), arrow_y + Pt(8),
                        color=color, width_pt=1.5)

    # 底部金句
    add_bottom_bar(slide, "架构哲学：零外部依赖 + 事件总线解耦 + 全栈 TypeScript = 生产级可演示系统",
                   highlight_words=["零外部依赖", "事件总线解耦", "全栈 TypeScript"])
