"""
第 5 页 · 需求与设计目标。

左半：5 个核心需求（图标 + 名称 + 简述）
右半：4 条设计原则（横向 4 卡）
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block, add_capsule


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=1, page_num=5)

    add_page_title(slide, "需求与设计目标", subtitle="5 个核心需求 + 4 条设计原则")

    # ---- 左半：5 个核心需求
    add_textbox(slide, left=Pt(80), top=Pt(180), width=Pt(400), height=Pt(30),
                text="5 个核心需求", font_size=15, bold=True, color=theme.PRIMARY)
    reqs = [
        ("画像", "6 维度对话式构建，随学随新",       "#FA8C16"),
        ("资源", "6 类资源多智能体协作生成",          "#52C41A"),
        ("路径", "AI 生成 + 12 预定义结构化路径",     "#1890FF"),
        ("辅导", "4 模式 + 追问链 + 画像注入",        "#722ED1"),
        ("评估", "真实进度同步 + 能力雷达 + 建议",    "#13C2C2"),
    ]
    for i, (title, desc, color) in enumerate(reqs):
        y = Pt(220) + i * Pt(70)
        add_card(slide, Pt(80), y, Pt(500), Pt(60), fill=theme.ACCENT_BG, border=None)
        add_color_block(slide, Pt(80), y, Pt(8), Pt(60), color)
        add_textbox(slide, Pt(100), y + Pt(8), Pt(80), Pt(24),
                    text=f"0{i+1}", font_size=18, bold=True, color=color)
        add_textbox(slide, Pt(180), y + Pt(8), Pt(380), Pt(24),
                    text=title, font_size=16, bold=True, color=theme.PRIMARY_DARK)
        add_textbox(slide, Pt(180), y + Pt(32), Pt(380), Pt(24),
                    text=desc, font_size=11, color=theme.TEXT_MUTED)

    # ---- 右半：4 条设计原则（2×2 卡片阵列）
    add_textbox(slide, left=Pt(640), top=Pt(180), width=Pt(530), height=Pt(30),
                text="4 条设计原则", font_size=15, bold=True, color="#FA8C16")
    principles = [
        ("个性化", "画像驱动所有下游智能体决策",     theme.PRIMARY),
        ("多智能体", "5 类智能体职责清晰、可组合",    "#FA8C16"),
        ("流式交互", "打字机效果 + 思考过程可视化",   "#722ED1"),
        ("数据闭环", "做题反馈回流画像、画像驱动推荐", "#13C2C2"),
    ]
    card_w = Pt(250)
    card_h = Pt(110)
    for i, (title, desc, color) in enumerate(principles):
        col = i % 2
        row = i // 2
        x = Pt(640) + col * (card_w + Pt(20))
        y = Pt(220) + row * (card_h + Pt(20))
        add_card(slide, x, y, card_w, card_h, fill=theme.WHITE, border=color, border_width=1.5)
        add_color_block(slide, x, y, card_w, Pt(6), color)
        add_textbox(slide, x + Pt(16), y + Pt(20), card_w - Pt(32), Pt(30),
                    text=title, font_size=18, bold=True, color=color)
        add_textbox(slide, x + Pt(16), y + Pt(55), card_w - Pt(32), Pt(45),
                    text=desc, font_size=11, color=theme.TEXT_MUTED)

    # ---- 底部 slogan
    add_rect(slide, Pt(80), Pt(630), Pt(1090), Pt(40), fill=theme.PRIMARY)
    add_textbox(slide, Pt(80), Pt(630), Pt(1090), Pt(40),
                text="让每个学生都拥有自己的 AI 学习智能体", font_size=16, bold=True,
                color=theme.WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)