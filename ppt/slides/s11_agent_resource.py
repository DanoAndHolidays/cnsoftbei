"""
第 11 页 · 资源生成智能体 · 学术商务版。
"""

from pptx.util import Pt

from components import theme
from components.layout import add_textbox, add_rect, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from slides.s10_agent_profile import _screenshot


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=3, page_num=11)

    color = theme.NAVY
    name = theme.AGENT_NAMES_CN["resource"]
    en = theme.AGENT_NAMES_EN["resource"]

    add_rect(slide, Pt(40), Pt(56), Pt(32), Pt(3), fill=color)
    add_textbox(slide, left=Pt(40), top=Pt(64), width=Pt(860), height=Pt(38),
                text=f"{name}  ({en})", font_size=26, bold=True, color=theme.DARK_TEXT)
    add_textbox(slide, left=Pt(40), top=Pt(104), width=Pt(860), height=Pt(20),
                text="系统的\"生产层\" — 基于画像定制 6 类资源 · 多智能体流水线协作 · SSE 实时状态回传",
                font_size=14, color=theme.TEXT_MUTED)

    _screenshot(slide, theme.SCREENSHOTS["resource1"], Pt(40), Pt(130), Pt(520), Pt(350))

    rx = Pt(590)
    rw = Pt(330)

    # 卡片 1
    add_card(slide, rx, Pt(130), rw, Pt(70), fill=theme.LIGHT_GRAY)
    add_textbox(slide, rx + Pt(14), Pt(136), rw - Pt(28), Pt(20),
                text="角色定位 — 系统的\"生产层\"", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(158), rw - Pt(28), Pt(36),
                text="基于画像为每个学习主题生成6类定制资源。Planner拆解任务→5类Worker并行执行。",
                font_size=12, color=theme.DARK_TEXT)

    # 卡片 2：6 种资源
    add_card(slide, rx, Pt(210), rw, Pt(142), fill=theme.LIGHT_GRAY)
    add_textbox(slide, rx + Pt(14), Pt(216), rw - Pt(28), Pt(20),
                text="6 种资源类型（2×3 网格）", font_size=14, bold=True, color=color)
    types = [
        ("文档", "结构化Markdown"), ("思维导图", "Mermaid格式"),
        ("测验", "自动判分+AI评阅"), ("阅读", "深度阅读+概念标注"),
        ("视频脚本", "分镜+时间戳"), ("代码案例", "可执行+注释"),
    ]
    for i, (cn_t, desc) in enumerate(types):
        col = i % 3
        row = i // 3
        x = rx + Pt(14) + col * Pt(108)
        y = Pt(244) + row * Pt(40)
        add_rect(slide, x, y + Pt(4), Pt(2), Pt(14), fill=color)
        add_textbox(slide, x + Pt(8), y, Pt(80), Pt(20),
                    text=cn_t, font_size=12, bold=True, color=theme.NAVY)
        add_textbox(slide, x + Pt(8), y + Pt(18), Pt(98), Pt(20),
                    text=desc, font_size=10, color=theme.TEXT_MUTED)

    # 卡片 3：协作流水线
    add_card(slide, rx, Pt(362), rw, Pt(120), fill=theme.LIGHT_GRAY)
    add_textbox(slide, rx + Pt(14), Pt(368), rw - Pt(28), Pt(20),
                text="多智能体协作流水线", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(390), rw - Pt(28), Pt(86),
                text="· Planner Agent解析目标→拆解子任务\n"
                     "· 5类Worker并行执行，各负责一类资源\n"
                     "· SSE实时回传每个Worker状态\n"
                     "· 失败自动重试3次（指数退避）\n"
                     "· AbortSignal一键取消\n"
                     "· 6 Worker并发完成 < 8s",
                font_size=11, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "资源智能体将\"千人千面\"落地为\"一人一案\"：每个学生获得画像驱动的专属学习资源",
                   highlight_words=["千人千面", "一人一案", "专属"])
