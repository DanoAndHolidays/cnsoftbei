"""
第 4 页 · 赛题对标。

上半：A3 赛题 5 项功能需求 → 本项目实现（5 行表格）
下半：A3 赛题 4 项非功能需求 → 本项目实现（4 行表格）
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_capsule


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=1, page_num=4)

    add_page_title(slide, "赛题对标", subtitle="A3 赛题 5 项功能 + 4 项非功能需求 → 本项目实现")

    # ---- 上半：功能需求表（5 行 × 3 列）
    add_textbox(slide, left=Pt(80), top=Pt(180), width=Pt(300), height=Pt(28),
                text="▶ 功能需求（5 项）", font_size=14, bold=True, color=theme.PRIMARY)
    func_rows = [
        ("1. 对话式学习画像（≥6 维度）", "Profile 页 / 6 维度 / 随学随新",            "#8"),
        ("2. 多智能体协同资源生成（≥5 类）", "Resource 页 / 6 类 / MultiAgentScheduler",  "#9 / #13"),
        ("3. 个性化学习路径规划",         "Path 页 / AI 自由生成 + 12 预定义结构化路径", "#10"),
        ("4. 智能辅导（可选加分）",       "Tutor 页 / 4 种解答模式 / 追问链",            "#11"),
        ("5. 学习效果评估（可选加分）",   "Assessment 页 / 真实进度 / 能力雷达",         "#12"),
    ]
    table_top = Pt(215)
    col_widths = [Pt(420), Pt(550), Pt(120)]
    col_x = [Pt(80), Pt(80) + col_widths[0], Pt(80) + col_widths[0] + col_widths[1]]
    # 表头
    headers = ["A3 赛题需求", "本项目实现", "PPT 页码"]
    for i, h in enumerate(headers):
        add_card(slide, col_x[i], table_top, col_widths[i], Pt(28), fill=theme.PRIMARY, border=None)
        add_textbox(slide, col_x[i] + Pt(8), table_top + Pt(4), col_widths[i] - Pt(16), Pt(20),
                    text=h, font_size=12, bold=True, color=theme.WHITE)
    # 数据行
    for ri, (req, impl, page) in enumerate(func_rows):
        y = table_top + Pt(28) + ri * Pt(28)
        bg = theme.ACCENT_BG if ri % 2 == 0 else theme.WHITE
        for i, val in enumerate([req, impl, page]):
            add_card(slide, col_x[i], y, col_widths[i], Pt(28), fill=bg, border=theme.DIVIDER, border_width=0.5)
            color = theme.PRIMARY if i == 2 else theme.TEXT
            bold = (i == 2)
            add_textbox(slide, col_x[i] + Pt(8), y + Pt(4), col_widths[i] - Pt(16), Pt(20),
                        text=val, font_size=11, color=color, bold=bold)

    # ---- 下半：非功能需求表（4 行）
    nfr_top = Pt(395)
    add_textbox(slide, left=Pt(80), top=nfr_top - Pt(28), width=Pt(300), height=Pt(28),
                text="▶ 非功能需求（4 项）", font_size=14, bold=True, color="#FA8C16")
    nfr_rows = [
        ("流式输出 / Markdown / 卡片化", "streamChatCompletion / MarkdownRenderer / Card", "#14"),
        ("开源协议标注",                  "第 19 页致谢中列出",                            "#19"),
        ("防幻觉机制",                    "题库预写 / AI 判分参考答案 / 引用追溯",          "#16"),
        ("响应时间 / 进度追踪",            "多智能体协作实时状态 + 流式增量",                "#9 / #14"),
    ]
    # 表头
    for i, h in enumerate(headers):
        add_card(slide, col_x[i], nfr_top, col_widths[i], Pt(28), fill="#FA8C16", border=None)
        add_textbox(slide, col_x[i] + Pt(8), nfr_top + Pt(4), col_widths[i] - Pt(16), Pt(20),
                    text=h, font_size=12, bold=True, color=theme.WHITE)
    for ri, (req, impl, page) in enumerate(nfr_rows):
        y = nfr_top + Pt(28) + ri * Pt(28)
        bg = theme.ACCENT_BG if ri % 2 == 0 else theme.WHITE
        for i, val in enumerate([req, impl, page]):
            add_card(slide, col_x[i], y, col_widths[i], Pt(28), fill=bg, border=theme.DIVIDER, border_width=0.5)
            color = theme.PRIMARY if i == 2 else theme.TEXT
            bold = (i == 2)
            add_textbox(slide, col_x[i] + Pt(8), y + Pt(4), col_widths[i] - Pt(16), Pt(20),
                        text=val, font_size=11, color=color, bold=bold)

    # ---- 底部 callout
    callout_y = Pt(605)
    add_card(slide, Pt(80), callout_y, Pt(1090), Pt(60), fill=theme.PRIMARY_LIGHT, border=theme.PRIMARY, border_width=1.0)
    add_textbox(slide, Pt(100), callout_y + Pt(8), Pt(1050), Pt(20),
                text="✓ A3 赛题 5+4 项需求全部实现，可选加分项（辅导/评估）均超额完成", font_size=13, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, Pt(100), callout_y + Pt(32), Pt(1050), Pt(20),
                text="✓ 5 个创新点详见第 17 页", font_size=11, color=theme.PRIMARY)