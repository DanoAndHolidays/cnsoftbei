"""
第 13 页 · 关键技术 1：多智能体协同框架。

左侧：核心类图（自绘）+ systemPrompt 节选
右侧：时序图（5 智能体协作流水线）
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from components.code_block import render_code_block
from components.flow_diagram import build_node, build_arrow


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=4, page_num=13)

    # 顶部徽章 + 标题
    add_color_block(slide, Pt(80), Pt(70), Pt(140), Pt(28), "#722ED1")
    add_textbox(slide, Pt(80), Pt(70), Pt(140), Pt(28),
                text="关键技术 01", font_size=14, bold=True, color=theme.WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left=Pt(232), top=Pt(70), width=Pt(900), height=Pt(38),
                text="多智能体协同框架（Multi-Agent Scheduler）",
                font_size=24, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, left=Pt(232), top=Pt(112), width=Pt(900), height=Pt(20),
                text="统一调度 5 类智能体，支持单智能体执行 / 流水线式协作",
                font_size=12, color=theme.TEXT_MUTED)

    # 左侧：核心类图 + 代码
    add_textbox(slide, Pt(80), Pt(170), Pt(560), Pt(24),
                text="▶ MultiAgentScheduler 核心类", font_size=14, bold=True, color="#722ED1")
    code = """class MultiAgentScheduler {
  private agents: Map<Role, Agent>;
  private eventBus: EventEmitter;

  registerAgent(role: Role, agent: Agent) {
    this.agents.set(role, agent);
    this.eventBus.on(`${role}:done`, this.onAgentDone);
  }

  async runPipeline(tasks: Task[]): Promise<Result[]> {
    const queue = [...tasks];
    const results = [];
    while (queue.length) {
      const task = queue.shift();
      const agent = this.agents.get(task.role);
      const result = await agent.execute(task, {
        signal: this.controller.signal,
      });
      results.push(result);
      this.eventBus.emit(`${task.role}:done`, result);
    }
    return results;
  }
}"""
    render_code_block(slide, Pt(80), Pt(200), Pt(560), Pt(280), code,
                      font_size=10, lang_label="TS")

    # 右侧：5 智能体流水线时序图
    add_textbox(slide, Pt(680), Pt(170), Pt(540), Pt(24),
                text="▶ 5 智能体协作流水线（资源生成为例）", font_size=14, bold=True, color="#722ED1")
    pipeline = ["planner", "document", "mindmap", "quiz", "reading", "video", "codeCase"]
    # 一行流程
    node_w = Pt(70)
    node_h = Pt(50)
    start_x = Pt(685)
    top = Pt(230)
    for i, name in enumerate(pipeline):
        x = start_x + i * (node_w + Pt(8))
        is_planner = (i == 0)
        fill = "#722ED1" if is_planner else theme.PRIMARY
        build_node(slide, x, top, node_w, node_h, name, fill=fill, font_size=10, radius=0.2)
        if i < len(pipeline) - 1:
            ax1 = x + node_w
            ax2 = x + node_w + Pt(8)
            ay = top + node_h // 2
            build_arrow(slide, ax1, ay, ax2, ay, color=theme.TEXT_MUTED, width_pt=1.0)

    # 时序说明
    add_card(slide, Pt(685), Pt(320), Pt(540), Pt(160), fill=theme.ACCENT_BG, border="#722ED1", border_width=1.0)
    add_textbox(slide, Pt(700), Pt(330), Pt(510), Pt(24),
                text="▶ 状态机", font_size=12, bold=True, color="#722ED1")
    add_textbox(slide, Pt(700), Pt(356), Pt(510), Pt(120),
                text="• planner 拆任务 → 派发给 6 个 worker\n"
                     "• worker 状态：pending → running → done/failed\n"
                     "• 失败自动重试 3 次\n"
                     "• 事件总线广播进度，UI 实时更新",
                font_size=11, color=theme.TEXT)

    # 底部 callout
    add_card(slide, Pt(80), Pt(610), Pt(1130), Pt(60),
             fill=theme.PRIMARY_LIGHT, border="#722ED1", border_width=1.0)
    add_textbox(slide, Pt(100), Pt(620), Pt(1090), Pt(20),
                text="✓ 关键设计：智能体之间解耦（仅通过事件总线通信），便于新增第 6、第 7 类智能体",
                font_size=12, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, Pt(100), Pt(642), Pt(1090), Pt(20),
                text="✓ 复用：ResourceGenerator 在 MultiAgentScheduler 之上封装，代码减少 40%",
                font_size=11, color=theme.PRIMARY)