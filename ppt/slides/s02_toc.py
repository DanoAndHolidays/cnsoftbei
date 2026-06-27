"""
第 2 页 · 目录 · 学术商务版。
5 个章节卡片 + 白色/浅灰交替底色。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from components import theme
from components.layout import add_textbox, add_rect, apply_chrome_v2
from components.shapes import add_card


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=1, page_num=2)

    # 蓝色短横线 + 标题
    add_rect(slide, left=Pt(40), top=Pt(58), width=Pt(36), height=Pt(3), fill=theme.NAVY)
    add_textbox(slide, left=Pt(40), top=Pt(68), width=Pt(860), height=Pt(42),
                text="目  录", font_size=36, bold=True, color=theme.NAVY)
    add_textbox(slide, left=Pt(40), top=Pt(114), width=Pt(860), height=Pt(22),
                text="CONTENTS  ·  21 页 / 5 章节  ·  导入 → 架构 → 智能体 → 技术 → 总结",
                font_size=14, color=theme.TEXT_MUTED)

    sections = [
        ("01", "项目导入与需求对标",
         "AI教育市场趋势 · 传统平台痛点 · 本项目差异化定位 · A3赛题对标", "P. 3 – 5"),
        ("02", "系统架构与技术选型",
         "四层轻量化架构（表现层/网关/智能体/数据层）· 技术选型与分析", "P. 6 – 8"),
        ("03", "五大核心智能体设计",
         "画像构建 · 资源生成 · 路径规划 · 智能辅导 · 效果评估", "P. 9 – 14"),
        ("04", "关键技术深挖与总结展望",
         "多智能体协同 · 流式交互 · 双向同步 · 系统评估 · 创新总结", "P. 15 – 20"),
        ("05", "致谢",
         "项目总结 · 开源组件致谢 · Q&A 交流", "P. 21"),
    ]

    row_top = Pt(160)
    row_h = Pt(72)
    row_gap = Pt(5)

    for i, (num, title, desc, pages) in enumerate(sections):
        y = row_top + i * (row_h + row_gap)
        bg = theme.LIGHT_GRAY if i % 2 == 0 else theme.WHITE
        add_card(slide, Pt(40), y, Pt(880), row_h, fill=bg, border=theme.BORDER, border_width=0.5)

        # 编号
        add_textbox(slide, left=Pt(54), top=y + Pt(6), width=Pt(80), height=Pt(60),
                    text=num, font_size=44, bold=True, color=theme.NAVY)

        add_rect(slide, left=Pt(144), top=y + Pt(14), width=Pt(2), height=row_h - Pt(28),
                 fill=theme.NAVY)

        # 章节名
        add_textbox(slide, left=Pt(160), top=y + Pt(8), width=Pt(520), height=Pt(30),
                    text=title, font_size=20, bold=True, color=theme.DARK_TEXT)

        # 描述
        add_textbox(slide, left=Pt(160), top=y + Pt(40), width=Pt(560), height=Pt(22),
                    text=desc, font_size=13, color=theme.TEXT_MUTED)

        # 页码
        add_textbox(slide, left=Pt(750), top=y + Pt(24), width=Pt(160), height=Pt(24),
                    text=pages, font_size=15, color=theme.NAVY, align=PP_ALIGN.RIGHT, bold=True)
