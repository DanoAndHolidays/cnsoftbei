"""
第 19 页 · 系统评估 · 学术商务版。
4 数据卡 + 5 测试结论 + 底部金句。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=4, page_num=19)

    add_page_title(slide, "系统评估与核心性能指标",
                   subtitle="576题库 · 92%判分一致率 · <2s首字延迟 · 0严重缺陷 · 6Worker并发<8s",
                   accent_color=theme.NAVY)

    # ── 4 个数据卡 ──
    add_textbox(slide, Pt(40), Pt(130), Pt(840), Pt(22),
                text="关键性能指标", font_size=16, bold=True, color=theme.NAVY)

    metrics = [
        ("576", "题库总题数", "12库×48题", theme.NAVY),
        ("92%", "AI判分一致率", "vs人工判分", theme.BLUE_MID),
        ("<2s", "流式首字延迟", "SSE首chunk", theme.NAVY),
        ("0", "已知严重缺陷", "build通过", theme.BLUE_MID),
    ]
    card_w = Pt(200)
    card_h = Pt(90)
    card_top = Pt(160)

    for i, (num, label, sub, color) in enumerate(metrics):
        x = Pt(40) + i * (card_w + Pt(14))
        add_card(slide, x, card_top, card_w, card_h, fill=theme.LIGHT_GRAY, border=color, border_width=1.5)
        add_rect(slide, x, card_top, card_w, Pt(4), fill=color)
        add_textbox(slide, x + Pt(16), card_top + Pt(12), card_w - Pt(32), Pt(36),
                    text=num, font_size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Pt(16), card_top + Pt(48), card_w - Pt(32), Pt(20),
                    text=label, font_size=14, bold=True, color=theme.DARK_TEXT, align=PP_ALIGN.CENTER)
        add_textbox(slide, x + Pt(16), card_top + Pt(68), card_w - Pt(32), Pt(18),
                    text=sub, font_size=11, color=theme.TEXT_MUTED, align=PP_ALIGN.CENTER)

    # ── 5 项测试结论 ──
    add_textbox(slide, Pt(40), Pt(268), Pt(840), Pt(22),
                text="5 项系统测试结论", font_size=16, bold=True, color=theme.NAVY)

    tests = [
        ("题库覆盖度", "12题库覆盖Python/Web/数据结构/网络/OS/数据库；题型6:3:1科学分布", theme.NAVY),
        ("AI判分一致性", "随机30题简答，AI vs 人工判分一致率92%；开放性题目偶有差异", theme.BLUE_MID),
        ("跨页面同步", "Practice↔Assessment同步验证通过；缓存命中100%；事件延迟<100ms", theme.NAVY),
        ("并发性能", "6Worker并发完成<8s(平均6.2s)；失败重试3次全恢复；SSE无丢失", theme.BLUE_MID),
        ("构建兼容性", "npm run build零错误；tsc类型检查通过；Win/Mac/Linux三平台兼容", theme.NAVY),
    ]
    test_h = Pt(36)
    test_top = Pt(296)

    for i, (title, desc, color) in enumerate(tests):
        y = test_top + i * (test_h + Pt(5))
        add_card(slide, Pt(40), y, Pt(880), test_h, fill=theme.WHITE, border=color, border_width=0.75)
        add_rect(slide, Pt(40), y, Pt(4), test_h, fill=color)
        add_textbox(slide, Pt(56), y + Pt(6), Pt(160), Pt(20),
                    text=title, font_size=13, bold=True, color=color)
        add_textbox(slide, Pt(230), y + Pt(8), Pt(680), Pt(22),
                    text=desc, font_size=12, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "系统评估结论：功能完备 + 性能达标 + 工程稳健 — 达到生产级交付标准",
                   highlight_words=["功能完备", "性能达标", "生产级"])
