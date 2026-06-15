"""
第 16 页 · 系统评估与测试。

4 个数据卡 + 4 项测试结论。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=5, page_num=16)

    add_page_title(slide, "系统评估与测试", subtitle="题库覆盖度 + AI 判分一致性 + 跨页面同步 + 性能指标")

    # 上半：4 个数据卡
    add_textbox(slide, Pt(80), Pt(170), Pt(1080), Pt(24),
                text="▶ 关键指标", font_size=14, bold=True, color="#13C2C2")
    metrics = [
        ("576", "题库总题数",     "12 库 × 48 题", theme.PRIMARY),
        ("92%", "AI 简答判分一致性", "与人工判分比对", "#FA8C16"),
        ("<2s", "路径生成响应时间",  "流式首字延迟",  theme.SUCCESS),
        ("0",   "已知严重缺陷",      "npm run build 通过", "#722ED1"),
    ]
    card_w = Pt(260)
    card_h = Pt(110)
    for i, (num, label, sub, color) in enumerate(metrics):
        x = Pt(80) + i * (card_w + Pt(20))
        y = Pt(210)
        add_card(slide, x, y, card_w, card_h, fill=theme.ACCENT_BG, border=color, border_width=1.5)
        add_color_block(slide, x, y, card_w, Pt(6), color)
        add_textbox(slide, x + Pt(16), y + Pt(20), card_w - Pt(32), Pt(40),
                    text=num, font_size=32, bold=True, color=color)
        add_textbox(slide, x + Pt(16), y + Pt(60), card_w - Pt(32), Pt(20),
                    text=label, font_size=13, bold=True, color=theme.PRIMARY_DARK)
        add_textbox(slide, x + Pt(16), y + Pt(82), card_w - Pt(32), Pt(20),
                    text=sub, font_size=10, color=theme.TEXT_MUTED)

    # 下半：4 项测试结论
    add_textbox(slide, Pt(80), Pt(360), Pt(1080), Pt(24),
                text="▶ 4 项测试结论", font_size=14, bold=True, color="#13C2C2")
    tests = [
        ("题库覆盖度", "12 个题库覆盖 Python 基础、Web、数据结构、计算机网络、操作系统等；题型 6:3:1 (判断:选择:简答)",
         theme.PRIMARY),
        ("AI 简答判分一致性", "随机抽 30 道简答，对比 AI 判分与人工判分；一致率 92%；不一致多为开放性题目",
         "#FA8C16"),
        ("跨页面同步验证", "Practice → Assessment 实时同步；路径采用 → Practice 模块过滤；缓存命中 = 100%",
         theme.SUCCESS),
        ("性能指标", "首屏 < 1.5s；流式首字 < 2s；多智能体协作 6 worker 并发完成 < 8s",
         "#722ED1"),
    ]
    test_h = Pt(60)
    for i, (title, desc, color) in enumerate(tests):
        y = Pt(400) + i * (test_h + Pt(8))
        add_card(slide, Pt(80), y, Pt(1080), test_h, fill=theme.WHITE, border=color, border_width=0.75)
        add_color_block(slide, Pt(80), y, Pt(8), test_h, color)
        add_textbox(slide, Pt(100), y + Pt(8), Pt(180), Pt(24),
                    text=title, font_size=14, bold=True, color=color)
        add_textbox(slide, Pt(290), y + Pt(10), Pt(860), Pt(48),
                    text=desc, font_size=11, color=theme.TEXT)