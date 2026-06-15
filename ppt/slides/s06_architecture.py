"""
第 6 页 · 总体架构（核心自绘图）。

4 层架构自上而下：表现层 → API 网关 → 多智能体框架 → 数据层。
每层一个圆角矩形 + 层名 + 包含组件。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from components.flow_diagram import build_node, build_arrow


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=2, page_num=6)

    add_page_title(slide, "总体架构", subtitle="4 层架构：表现层 → API 网关 → 多智能体框架 → 数据层")

    # ---- 4 层架构图
    layers = [
        ("表现层 (Presentation)", "React 19 + TypeScript + Vite + Ant Design 6",                  theme.PRIMARY),
        ("API 网关 (Gateway)",    "Vite dev proxy: /anthropic → api.minimaxi.com",                 "#9254DE"),
        ("多智能体框架 (Multi-Agent)", "MultiAgentScheduler  ·  5 类智能体  ·  事件总线",          "#FA8C16"),
        ("数据层 (Data)",         "localStorage  ·  题库 JSON（12 库 576 题）",                   theme.SUCCESS),
    ]
    layer_top = Pt(190)
    layer_h = Pt(95)
    layer_gap = Pt(30)
    for i, (name, content, color) in enumerate(layers):
        y = layer_top + i * (layer_h + layer_gap)
        # 左侧色条
        add_rect(slide, Pt(140), y, Pt(12), layer_h, fill=color)
        # 主体卡片
        add_card(slide, Pt(160), y, Pt(900), layer_h, fill=theme.ACCENT_BG, border=color, border_width=1.5)
        add_textbox(slide, Pt(180), y + Pt(15), Pt(380), Pt(30),
                    text=name, font_size=18, bold=True, color=color)
        add_textbox(slide, Pt(180), y + Pt(48), Pt(860), Pt(30),
                    text=content, font_size=13, color=theme.TEXT)
        # 右侧层编号大数字
        add_textbox(slide, Pt(1010), y + Pt(15), Pt(120), layer_h - Pt(30),
                    text=f"L{i+1}", font_size=36, bold=True, color=color, align=PP_ALIGN.RIGHT)
        # 层间箭头（除最后一层）
        if i < len(layers) - 1:
            arrow_y = y + layer_h + Pt(8)
            build_arrow(slide, Pt(610), arrow_y, Pt(610), arrow_y + Pt(14), color=color, width_pt=2.0)

    # ---- 底部右侧 callout
    add_card(slide, Pt(700), Pt(640), Pt(500), Pt(50),
             fill=theme.PRIMARY_LIGHT, border=theme.PRIMARY, border_width=1.0)
    add_textbox(slide, Pt(720), Pt(648), Pt(460), Pt(20),
                text="✓ 全栈可本地运行：npm run dev 一行启动", font_size=12, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, Pt(720), Pt(668), Pt(460), Pt(20),
                text="✓ 无外部数据库依赖（localStorage 持久化）", font_size=12, color=theme.PRIMARY)
