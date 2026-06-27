"""
第 21 页 · 致谢 · 学术商务版。
深蓝全幅 + 白字 + 项目信息 + Q&A。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from components import theme
from components.layout import add_textbox, add_rect, hex_to_rgb


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # 深蓝全幅
    add_rect(slide, left=Pt(0), top=Pt(0),
             width=theme.SLIDE_WIDTH, height=theme.SLIDE_HEIGHT,
             fill=theme.NAVY)

    # 顶部红色细线
    add_rect(slide, left=Pt(0), top=Pt(0),
             width=theme.SLIDE_WIDTH, height=Pt(2.5),
             fill=theme.RED)

    # 大字
    add_textbox(slide, left=Pt(80), top=Pt(60), width=Pt(800), height=Pt(90),
                text="感谢聆听", font_size=64, bold=True,
                color=theme.WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, left=Pt(80), top=Pt(156), width=Pt(800), height=Pt(24),
                text="THANK YOU FOR YOUR ATTENTION",
                font_size=14, color="#8899BB", align=PP_ALIGN.CENTER)

    add_rect(slide, left=Pt(460), top=Pt(190), width=Pt(40), height=Pt(2.5),
             fill=theme.RED)

    # 项目概述
    add_textbox(slide, left=Pt(80), top=Pt(222), width=Pt(800), height=Pt(22),
                text="项目概述", font_size=18, bold=True, color=theme.WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, left=Pt(100), top=Pt(254), width=Pt(760), height=Pt(60),
                text="赛题：第十五届中国软件杯·A3赛题  |  作品：学习智能体系统\n"
                     "技术栈：React 19 + TypeScript + Vite 5 + Ant Design 6 + Recharts\n"
                     "核心数据：576题库 · 92%判分一致率 · <2s首字延迟 · 0严重缺陷",
                font_size=13, color="#AABBCC", align=PP_ALIGN.CENTER)

    # 致谢
    add_textbox(slide, left=Pt(80), top=Pt(340), width=Pt(800), height=Pt(22),
                text="致谢", font_size=18, bold=True, color=theme.WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, left=Pt(100), top=Pt(370), width=Pt(760), height=Pt(50),
                text="感谢评委老师的辛勤评审  |  感谢指导老师的悉心指导\n"
                     "本作品使用：React·TypeScript·Vite·Ant Design·Recharts·python-pptx·MiniMax API",
                font_size=12, color="#8899BB", align=PP_ALIGN.CENTER)

    # Q&A
    add_rect(slide, left=Pt(380), top=Pt(440), width=Pt(200), height=Pt(40),
             fill=theme.RED)
    add_textbox(slide, left=Pt(380), top=Pt(442), width=Pt(200), height=Pt(36),
                text="Q & A", font_size=20, bold=True,
                color=theme.WHITE, align=PP_ALIGN.CENTER)

    # 底部信息
    add_textbox(slide, left=Pt(80), top=Pt(500), width=Pt(800), height=Pt(22),
                text=f"{theme.COVER_INFO['team_name']}  ·  {theme.COVER_INFO['school']}",
                font_size=15, bold=True, color=theme.WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, left=Pt(80), top=Pt(522), width=Pt(800), height=Pt(16),
                text=f"{theme.COVER_INFO['date']}  ·  第十五届中国软件杯·A3赛题  ·  {theme.COVER_INFO['presenter']} 汇报",
                font_size=11, color="#667799", align=PP_ALIGN.CENTER)

    # 装饰圆
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Pt(-60), Pt(380), Pt(180), Pt(180))
    shp.fill.solid()
    shp.fill.fore_color.rgb = hex_to_rgb("#2A4A9A")
    shp.fill.transparency = 0.75
    shp.line.fill.background()

    shp2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Pt(840), Pt(380), Pt(180), Pt(180))
    shp2.fill.solid()
    shp2.fill.fore_color.rgb = hex_to_rgb("#2A4A9A")
    shp2.fill.transparency = 0.75
    shp2.line.fill.background()
