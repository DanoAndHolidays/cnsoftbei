"""
第 20 页 · 创新总结与未来展望 · 学术商务版。
上：5 项创新  |  下：3 大方向  +  底部金句。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=4, page_num=20)

    add_page_title(slide, "创新价值沉淀与未来演进方向",
                   subtitle="架构/数据/交互/工程 4维度 5项创新 + 3大未来方向",
                   accent_color=theme.NAVY)

    # ── 5 项核心创新 ──
    add_textbox(slide, Pt(40), Pt(130), Pt(880), Pt(22),
                text="5 项核心创新沉淀", font_size=16, bold=True, color=theme.NAVY)

    innovations = [
        ("01", "5智能体\n协同框架", "MultiAgentScheduler调度\nEventEmitter解耦\n代码量减少40%"),
        ("02", "6维动态\n画像系统", "对话式+做题回流双更新\nlocalStorage跨页面共享\n下游prompt注入"),
        ("03", "结构化\n路径节点", "绑定题库/模块ID\n80%阈值自动标记\n12预定义+AI生成"),
        ("04", "流式思考\n可视化", "SSE打字机效果\n<thinking>块可折叠\nAbortSignal中断"),
        ("05", "Tutor\n5项优化", "画像注入+缓存去重\n追问链+点踩重生\n4模式+生产级稳定"),
    ]

    cap_w = Pt(168)
    cap_h = Pt(174)
    cap_top = Pt(158)
    cap_gap = Pt(8)

    for i, (num, title, body) in enumerate(innovations):
        x = Pt(40) + i * (cap_w + cap_gap)
        add_card(slide, x, cap_top, cap_w, cap_h, fill=theme.WHITE, border=theme.NAVY, border_width=1.5)
        add_rect(slide, x, cap_top, cap_w, Pt(4), fill=theme.NAVY)
        add_textbox(slide, x + Pt(10), cap_top + Pt(10), cap_w - Pt(20), Pt(22),
                    text=num, font_size=18, bold=True, color=theme.NAVY)
        add_textbox(slide, x + Pt(10), cap_top + Pt(34), cap_w - Pt(20), Pt(40),
                    text=title, font_size=13, bold=True, color=theme.DARK_TEXT)
        add_textbox(slide, x + Pt(10), cap_top + Pt(76), cap_w - Pt(20), cap_h - Pt(84),
                    text=body, font_size=11, color=theme.TEXT_MUTED)

    # ── 3 大未来方向 ──
    add_textbox(slide, Pt(40), Pt(348), Pt(880), Pt(22),
                text="3 大未来演进方向", font_size=16, bold=True, color=theme.RED)

    futures = [
        ("01", "多模态扩展", "拍照搜题+语音问答",
         "接入图像识别与语音交互大模型；拍照→OCR→解答；语音→ASR→TTS；手写批改；AI生成教学视频"),
        ("02", "知识图谱构建", "标签匹配→逻辑推理",
         "构建学科知识图谱(概念→前置→后继)；路径推荐基于图遍历；跨学科关联；知识点掌握度传播"),
        ("03", "跨用户协作网络", "个体→社交化协同",
         "学习小组+目标+进度看板；同伴互评AI辅助仲裁；错题共享+讨论；排行榜与成就系统"),
    ]

    card_w = Pt(280)
    card_h = Pt(112)
    card_top = Pt(376)

    for i, (num, title, subtitle, body) in enumerate(futures):
        x = Pt(40) + i * (card_w + Pt(16))
        add_card(slide, x, card_top, card_w, card_h, fill=theme.LIGHT_GRAY, border=theme.NAVY, border_width=1.5)
        add_rect(slide, x, card_top, Pt(4), card_h, fill=theme.NAVY)
        add_textbox(slide, x + Pt(14), card_top + Pt(6), Pt(40), Pt(22),
                    text=num, font_size=18, bold=True, color=theme.NAVY)
        add_textbox(slide, x + Pt(52), card_top + Pt(6), card_w - Pt(66), Pt(22),
                    text=title, font_size=15, bold=True, color=theme.DARK_TEXT)
        add_textbox(slide, x + Pt(14), card_top + Pt(30), card_w - Pt(28), Pt(18),
                    text=subtitle, font_size=12, bold=True, color=theme.RED)
        add_textbox(slide, x + Pt(14), card_top + Pt(50), card_w - Pt(28), card_h - Pt(56),
                    text=body, font_size=11, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "愿景：AI+教育的终局不是替代教师，而是让每个学生都拥有最懂自己的专属学习智能体",
                   highlight_words=["专属学习智能体"])
