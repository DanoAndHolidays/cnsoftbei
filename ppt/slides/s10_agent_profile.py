"""
第 10 页 · 画像构建智能体 · 学术商务版。
左截图 + 右 3 卡片 + 底部金句。
"""

import os
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from PIL import Image


def _screenshot(slide, img_rel_path: str, left, top, width, height):
    ppt_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    img_path = os.path.abspath(os.path.join(ppt_dir, theme.SCREENSHOT_DIR, img_rel_path))
    if not os.path.exists(img_path):
        add_card(slide, left, top, width, height, fill=theme.LIGHT_GRAY)
        add_textbox(slide, left, top, width, height,
                    text=f"[截图缺失]", font_size=11, color=theme.TEXT_SUBTLE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        return
    img = Image.open(img_path)
    iw, ih = img.size
    ratio = min(width / iw, height / ih)
    w = int(iw * ratio)
    h = int(ih * ratio)
    x = left + (width - w) // 2
    y = top + (height - h) // 2
    slide.shapes.add_picture(img_path, x, y, width=w, height=h)


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=3, page_num=10)

    color = theme.NAVY
    name = theme.AGENT_NAMES_CN["profile"]
    en = theme.AGENT_NAMES_EN["profile"]

    # 标题
    add_rect(slide, Pt(40), Pt(56), Pt(32), Pt(3), fill=color)
    add_textbox(slide, left=Pt(40), top=Pt(64), width=Pt(860), height=Pt(38),
                text=f"{name}  ({en})", font_size=26, bold=True, color=theme.DARK_TEXT)
    add_textbox(slide, left=Pt(40), top=Pt(104), width=Pt(860), height=Pt(20),
                text="系统的\"感知层\" — 对话式 6 维画像构建 · 随学随新动态更新 · 跨页面持久化共享",
                font_size=14, color=theme.TEXT_MUTED)

    # 左侧截图
    _screenshot(slide, theme.SCREENSHOTS["profile"], Pt(40), Pt(130), Pt(520), Pt(350))

    # 右侧
    rx = Pt(590)
    rw = Pt(330)

    # 卡片 1：角色定位
    add_card(slide, rx, Pt(130), rw, Pt(72), fill=theme.LIGHT_GRAY)
    add_textbox(slide, rx + Pt(14), Pt(136), rw - Pt(28), Pt(20),
                text="角色定位 — 系统的\"感知层\"", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(158), rw - Pt(28), Pt(38),
                text="通过对话理解学习者，构建6维动态画像；画像随学习过程持续更新，注入所有下游智能体。",
                font_size=12, color=theme.DARK_TEXT)

    # 卡片 2：6 维画像
    add_card(slide, rx, Pt(212), rw, Pt(154), fill=theme.LIGHT_GRAY)
    add_textbox(slide, rx + Pt(14), Pt(218), rw - Pt(28), Pt(20),
                text="6 维画像结构", font_size=14, bold=True, color=color)
    dims = [
        ("知识基础", "已掌握知识的广度与深度"),
        ("认知风格", "视觉/文字/动手偏好"),
        ("易错偏好", "高频错误与薄弱区"),
        ("学习节奏", "最佳学习时长与速度"),
        ("兴趣方向", "偏好的学科与技术方向"),
        ("学习习惯", "时间段/专注度/持续性"),
    ]
    for i, (title, desc) in enumerate(dims):
        col = i % 3
        row = i // 3
        x = rx + Pt(14) + col * Pt(108)
        y = Pt(244) + row * Pt(40)
        add_rect(slide, x, y + Pt(4), Pt(2), Pt(12), fill=color)
        add_textbox(slide, x + Pt(8), y, Pt(82), Pt(18),
                    text=title, font_size=12, bold=True, color=color)
        add_textbox(slide, x + Pt(8), y + Pt(18), Pt(98), Pt(20),
                    text=desc, font_size=10, color=theme.TEXT_MUTED)

    # 卡片 3：关键特性
    add_card(slide, rx, Pt(376), rw, Pt(106), fill=theme.LIGHT_GRAY)
    add_textbox(slide, rx + Pt(14), Pt(382), rw - Pt(28), Pt(20),
                text="关键特性：随学随新 + 跨智能体共享", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(404), rw - Pt(28), Pt(72),
                text="· 对话式构建：自然语言输入\n"
                     "· localStorage跨页面持久化\n"
                     "· 做题反馈回流自动更新画像\n"
                     "· System Prompt注入所有下游智能体\n"
                     "· 6维雷达可视化展示",
                font_size=11, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "画像智能体是整个系统的\"大脑\"：从感知到决策，驱动所有下游智能体的个性化行为",
                   highlight_words=["大脑", "个性化行为"])
