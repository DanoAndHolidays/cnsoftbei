"""
第 19 页 · 致谢 / Q&A。

简洁封底：Thanks + 队员 + 联系方式 + Q&A 提示。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect
from components.shapes import add_card, add_color_block


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=5, page_num=19)

    # 不调用 add_page_title，用大字 Thanks 居中

    # 主背景色（浅蓝灰）
    add_rect(slide, Pt(0), Pt(0), theme.SLIDE_WIDTH, theme.SLIDE_HEIGHT, fill=theme.ACCENT_BG)

    # 重新叠加左侧装饰条（因为上面填充了背景）
    add_rect(slide, Pt(0), Pt(50), Pt(4), theme.SLIDE_HEIGHT - Pt(100), fill="#13C2C2")

    # 巨大 Thanks
    add_textbox(slide, Pt(0), Pt(180), theme.SLIDE_WIDTH, Pt(160),
                text="Thanks", font_size=120, bold=True, color=theme.PRIMARY,
                align=PP_ALIGN.CENTER)
    add_textbox(slide, Pt(0), Pt(360), theme.SLIDE_WIDTH, Pt(40),
                text="感谢聆听 · 欢迎提问", font_size=24, color=theme.PRIMARY_DARK,
                align=PP_ALIGN.CENTER)

    # 4px 橙色短分割线
    add_rect(slide, left=Pt(620), top=Pt(420), width=Pt(80), height=Pt(4), fill="#FA8C16")

    # 队员信息
    add_textbox(slide, Pt(0), Pt(460), theme.SLIDE_WIDTH, Pt(28),
                text=f"队伍：{theme.COVER_INFO['team_name']}    ·    学校：{theme.COVER_INFO['school']}",
                font_size=14, color=theme.TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, Pt(0), Pt(490), theme.SLIDE_WIDTH, Pt(28),
                text=f"汇报人：{theme.COVER_INFO['presenter']}    ·    指导老师：{theme.COVER_INFO['advisor']}",
                font_size=14, color=theme.TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, Pt(0), Pt(520), theme.SLIDE_WIDTH, Pt(28),
                text=f"联系方式：{theme.COVER_INFO['contact']}",
                font_size=14, color=theme.TEXT_MUTED, align=PP_ALIGN.CENTER)

    # 底部 Q&A 提示
    add_textbox(slide, Pt(0), Pt(620), theme.SLIDE_WIDTH, Pt(40),
                text="Q & A", font_size=36, bold=True, color=theme.PRIMARY_LIGHT,
                align=PP_ALIGN.CENTER)

    # 开源致谢
    add_textbox(slide, Pt(0), Pt(680), theme.SLIDE_WIDTH, Pt(20),
                text="本项目使用 React / TypeScript / Vite / Ant Design / Recharts / python-pptx / matplotlib 等开源组件，特此致谢",
                font_size=9, color=theme.TEXT_SUBTLE, align=PP_ALIGN.CENTER)