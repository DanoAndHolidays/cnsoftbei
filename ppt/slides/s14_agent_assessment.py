"""
第 14 页 · 效果评估智能体 · 学术商务版。
"""

from pptx.util import Pt

from components import theme
from components.layout import add_textbox, add_rect, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card
from components.assessment_mock import render_assessment_mock


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=3, page_num=14)

    color = theme.NAVY
    name = theme.AGENT_NAMES_CN["assessment"]
    en = theme.AGENT_NAMES_EN["assessment"]

    add_rect(slide, Pt(40), Pt(56), Pt(32), Pt(3), fill=color)
    add_textbox(slide, left=Pt(40), top=Pt(64), width=Pt(860), height=Pt(38),
                text=f"{name}  ({en})", font_size=26, bold=True, color=theme.DARK_TEXT)
    add_textbox(slide, left=Pt(40), top=Pt(104), width=Pt(860), height=Pt(20),
                text="系统的\"反馈层\" — 真实练习数据驱动 · 6 维能力雷达 · 智能调整建议 · 以评促学",
                font_size=14, color=theme.TEXT_MUTED)

    # 左侧 mockup
    render_assessment_mock(slide, Pt(40), Pt(130), Pt(520), Pt(350))

    # 右侧
    rx = Pt(590)
    rw = Pt(330)

    # 卡片 1
    add_card(slide, rx, Pt(130), rw, Pt(130), fill=theme.LIGHT_GRAY, border=color, border_width=1.0)
    add_textbox(slide, rx + Pt(14), Pt(136), rw - Pt(28), Pt(20),
                text="真实进度同步 — \"不造假\"原则", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(158), rw - Pt(28), Pt(96),
                text="· 直接读取practiceState（Practice页写入）\n"
                     "· 监听practiceStateUpdated事件实时刷新\n"
                     "· 4个统计卡+模块进度卡全基于真实数据\n"
                     "· 无练习记录→显示\"去练习\"引导提示\n"
                     "· 画像数据同步读取，雷达图反映最新状态",
                font_size=11, color=theme.DARK_TEXT)

    # 卡片 2
    add_card(slide, rx, Pt(272), rw, Pt(108), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, rx + Pt(14), Pt(278), rw - Pt(28), Pt(20),
                text="6 维能力雷达 + 智能建议", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(300), rw - Pt(28), Pt(74),
                text="【6维雷达Recharts RadarChart】\n"
                     "知识基础·认知风格·易错偏好\n"
                     "学习节奏·兴趣方向·学习习惯\n\n"
                     "【4类智能调整建议】\n"
                     "专项练习·节奏建议·路径推荐·整体策略",
                font_size=11, color=theme.DARK_TEXT)

    # 卡片 3
    add_card(slide, rx, Pt(392), rw, Pt(90), fill=theme.LIGHT_GRAY, border=color, border_width=1.0)
    add_textbox(slide, rx + Pt(14), Pt(398), rw - Pt(28), Pt(20),
                text="以评促学 — 学习闭环最后一环", font_size=14, bold=True, color=color)
    add_textbox(slide, rx + Pt(14), Pt(420), rw - Pt(28), Pt(56),
                text="评估不是终点，而是新起点：\n"
                     "暴露弱项→专项练习→节奏调整→路径推荐\n"
                     "形成\"练习→评估→建议→再练习\"持续改进循环",
                font_size=11, color=theme.DARK_TEXT)

    add_bottom_bar(slide, "评估智能体让数据\"说话\"：从练习到雷达再到建议，形成完整的\"以评促学\"学习闭环",
                   highlight_words=["以评促学", "数据\"说话\""])
