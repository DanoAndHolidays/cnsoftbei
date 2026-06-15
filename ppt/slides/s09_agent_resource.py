"""
第 9 页 · 资源生成智能体（Resource Agent）🟢
"""

import os
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from slides.s08_agent_profile import _screenshot


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=3, page_num=9)

    color = theme.AGENT_COLORS["resource"]
    name = theme.AGENT_NAMES_CN["resource"]
    en = theme.AGENT_NAMES_EN["resource"]
    emoji = theme.AGENT_EMOJI["resource"]

    add_color_block(slide, Pt(80), Pt(70), Pt(10), Pt(28), color)
    add_textbox(slide, left=Pt(100), top=Pt(70), width=Pt(700), height=Pt(38),
                text=f"{name} {emoji}", font_size=24, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, left=Pt(100), top=Pt(110), width=Pt(600), height=Pt(20),
                text=en, font_size=12, color=theme.TEXT_SUBTLE)

    _screenshot(slide, theme.SCREENSHOTS["resource1"], Pt(80), Pt(160), Pt(560), Pt(440))

    right_x = Pt(680)
    right_w = Pt(540)

    add_card(slide, right_x, Pt(160), right_w, Pt(110), fill=theme.ACCENT_BG, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(170), right_w - Pt(32), Pt(24),
                text="▶ 角色定位", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(196), right_w - Pt(32), Pt(70),
                text="基于画像为每个学习主题生成 6 类定制资源，多智能体协作流水线。",
                font_size=12, color=theme.TEXT)

    add_card(slide, right_x, Pt(280), right_w, Pt(160), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(290), right_w - Pt(32), Pt(24),
                text="▶ 6 种资源类型", font_size=13, bold=True, color=color)
    types = [("document", "文档"), ("mindmap", "思维导图"), ("quiz", "测验"),
             ("reading", "阅读"), ("video", "视频脚本"), ("codeCase", "代码案例")]
    for i, (en_t, cn_t) in enumerate(types):
        col = i % 2
        row = i // 2
        x = right_x + Pt(16) + col * Pt(260)
        y = Pt(320) + row * Pt(32)
        add_rect(slide, x, y + Pt(4), Pt(8), Pt(8), fill=color)
        add_textbox(slide, x + Pt(16), y, Pt(240), Pt(20),
                    text=f"{cn_t} ({en_t})", font_size=12, color=theme.TEXT)

    add_card(slide, right_x, Pt(450), right_w, Pt(150), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(460), right_w - Pt(32), Pt(24),
                text="▶ 多智能体协作（实时状态）", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(488), right_w - Pt(32), Pt(110),
                text="• planner 拆任务 → 5 类 worker 并行\n"
                     "• SSE 流式回传各 worker 状态（pending/running/done）\n"
                     "• 失败自动重试，最多 3 次\n"
                     "• 全部完成前允许用户中断（AbortSignal）",
                font_size=12, color=theme.TEXT)
