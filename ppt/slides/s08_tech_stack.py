"""
第 8 页 · 技术选型 · 学术商务版。
3 列布局 + 底部金句。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=2, page_num=8)

    add_page_title(slide, "前沿与稳健并重的技术选型",
                   subtitle="React 19 + SSE 流式 + Recharts 三大技术支撑")

    cols = [
        ("前端框架与工程化", theme.NAVY, [
            ("React 19", "最新稳定版·Concurrent·批处理"),
            ("TypeScript", "类型安全·IDE智能提示"),
            ("Vite 5", "极速HMR·冷启动<1s"),
            ("Ant Design 6", "企业级UI·Statistic/Progress"),
            ("PageCache", "手动路由·useRef状态缓存"),
        ]),
        ("AI 交互与协议层", theme.BLUE_MID, [
            ("大模型 API", "兼容Anthropic协议"),
            ("SSE 流式传输", "ReadableStream打字机效果"),
            ("AbortSignal", "中途取消请求+状态恢复"),
            ("Prompt 工程", "5类系统prompt·画像注入"),
            ("容错机制", "3次重试·3min超时·指数退避"),
        ]),
        ("数据可视化与存储", theme.NAVY, [
            ("Recharts 3", "React原生·RadarChart雷达"),
            ("雷达图", "6维度能力评估可视化"),
            ("统计卡片", "Antd Statistic数据展示"),
            ("时间线", "Antd Timeline智能建议"),
            ("localStorage", "零数据库·画像/进度/路径"),
        ]),
    ]

    col_w = Pt(280)
    col_top = Pt(130)

    for ci, (col_name, color, items) in enumerate(cols):
        x = Pt(40) + ci * (col_w + Pt(20))

        # 列标题
        add_card(slide, x, col_top, col_w, Pt(36), fill=color, border=None)
        add_textbox(slide, x, col_top, col_w, Pt(36),
                    text=col_name, font_size=15, bold=True, color=theme.WHITE,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 项目
        for i, (name, desc) in enumerate(items):
            y = col_top + Pt(48) + i * Pt(58)
            add_card(slide, x, y, col_w, Pt(50), fill=theme.LIGHT_GRAY, border=color, border_width=0.5)
            add_rect(slide, x, y, Pt(3), Pt(50), fill=color)
            add_textbox(slide, x + Pt(12), y + Pt(4), col_w - Pt(24), Pt(22),
                        text=name, font_size=13, bold=True, color=theme.NAVY)
            add_textbox(slide, x + Pt(12), y + Pt(28), col_w - Pt(24), Pt(18),
                        text=desc, font_size=11, color=theme.TEXT_MUTED)

    add_bottom_bar(slide, "选型哲学：主流成熟保稳定 + 前沿技术提体验 + 轻量存储降依赖 = 生产级可演示系统",
                   highlight_words=["主流成熟", "前沿技术", "轻量存储"])
