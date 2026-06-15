"""
第 18 页 · 总结与未来展望。

左：已完成清单（5 智能体 + 7 页面 + 12 题库）
右：3 个未来方向（多模态 / 知识图谱 / 协作学习）
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=5, page_num=18)

    add_page_title(slide, "总结与未来展望", subtitle="已完成 + 待拓展")

    # 左半：已完成清单
    add_textbox(slide, Pt(80), Pt(170), Pt(540), Pt(28),
                text="▶ 已完成", font_size=15, bold=True, color=theme.SUCCESS)
    done = [
        ("5", "智能体",   "画像 / 资源 / 路径 / 辅导 / 评估"),
        ("7", "核心页面", "Home / Profile / Resources / Path / Practice / Tutor / Assessment"),
        ("12", "题库",     "576 题，覆盖 Python / Web / 数据结构 / 计算机网络 / OS / 算法 等"),
        ("19", "PPT 页",   "本汇报稿"),
        ("5", "工程优化", "Tutor 缓存/追问链/点踩/取消/画像注入"),
    ]
    for i, (num, label, desc) in enumerate(done):
        y = Pt(220) + i * Pt(70)
        add_card(slide, Pt(80), y, Pt(540), Pt(60), fill=theme.ACCENT_BG, border=theme.SUCCESS, border_width=0.75)
        add_textbox(slide, Pt(96), y + Pt(8), Pt(60), Pt(40),
                    text=num, font_size=24, bold=True, color=theme.SUCCESS)
        add_textbox(slide, Pt(160), y + Pt(10), Pt(140), Pt(20),
                    text=label, font_size=14, bold=True, color=theme.PRIMARY_DARK)
        add_textbox(slide, Pt(160), y + Pt(32), Pt(440), Pt(24),
                    text=desc, font_size=10, color=theme.TEXT_MUTED)

    # 右半：未来方向
    add_textbox(slide, Pt(660), Pt(170), Pt(540), Pt(28),
                text="▶ 未来方向", font_size=15, bold=True, color="#FA8C16")
    futures = [
        ("多模态扩展",  "接入图像 / 语音识别：拍照搜题、语音问答",            "#FA8C16"),
        ("知识图谱",    "构建学科知识图谱，路径推荐更智能",                  theme.PRIMARY),
        ("跨用户协作",  "学习小组 / 同伴互评 / 错题共享",                    "#722ED1"),
    ]
    for i, (title, desc, color) in enumerate(futures):
        y = Pt(220) + i * Pt(110)
        add_card(slide, Pt(660), y, Pt(540), Pt(100), fill=theme.WHITE, border=color, border_width=1.5)
        add_color_block(slide, Pt(660), y, Pt(8), Pt(100), color)
        add_textbox(slide, Pt(680), y + Pt(12), Pt(200), Pt(30),
                    text=f"0{i+1}  {title}", font_size=16, bold=True, color=color)
        add_textbox(slide, Pt(680), y + Pt(48), Pt(500), Pt(48),
                    text=desc, font_size=12, color=theme.TEXT)

    # 底部 slogan
    add_rect(slide, Pt(80), Pt(620), Pt(1130), Pt(50), fill=theme.PRIMARY)
    add_textbox(slide, Pt(80), Pt(620), Pt(1130), Pt(50),
                text="我们相信：AI + 教育 = 每个学生都有专属的学习智能体",
                font_size=18, bold=True, color=theme.WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)