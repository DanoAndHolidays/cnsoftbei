"""
第 3 页 · 项目背景。

左：AI 教育趋势（3 个数据点）
中：传统平台 3 大痛点（漏斗示意）
右：本项目定位（差异化 4 点）
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from components.flow_diagram import build_node, build_arrow


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=1, page_num=3)

    add_page_title(slide, "项目背景", subtitle="AI 教育趋势 + 传统平台痛点 + 本项目定位")

    # ---- 左栏：AI 教育趋势（3 个数据卡）
    left_x = Pt(80)
    left_w = Pt(360)
    col_top = Pt(200)
    add_textbox(slide, left_x, col_top, left_w, Pt(30),
                "AI 教育市场趋势",
                font_size=15, bold=True, color=theme.PRIMARY_DARK)
    trends = [
        ("60%+", "学生认为 AI 个性化辅导效果优于传统课堂"),
        ("3 倍", "近 2 年自适应学习平台用户增长"),
        ("80%", "高校将引入 AI 教学辅助系统（2026 预测）"),
    ]
    for i, (num, desc) in enumerate(trends):
        y = col_top + Pt(45) + i * Pt(90)
        add_card(slide, left_x, y, left_w, Pt(80), fill=theme.ACCENT_BG)
        add_textbox(slide, left_x + Pt(16), y + Pt(8), left_w - Pt(32), Pt(34),
                    num,
                    font_size=24, bold=True, color=theme.PRIMARY)
        add_textbox(slide, left_x + Pt(16), y + Pt(44), left_w - Pt(32), Pt(32),
                    desc,
                    font_size=11, color=theme.TEXT_MUTED)

    # ---- 中栏：传统平台 3 大痛点（漏斗：3 个倒梯形）
    mid_x = Pt(480)
    mid_w = Pt(360)
    add_textbox(slide, mid_x, col_top, mid_w, Pt(30),
                "传统学习平台 3 大痛点",
                font_size=15, bold=True, color=theme.PRIMARY_DARK)
    pains = [
        ("资源繁杂", "题库、视频、文档分散", theme.ERROR),
        ("节奏统一", "全班同一进度，难个性化", theme.WARNING),
        ("反馈滞后", "错题几天后才讲评", "#722ED1"),
    ]
    pain_top = col_top + Pt(50)
    for i, (title, desc, color) in enumerate(pains):
        y = pain_top + i * Pt(90)
        w = mid_w - i * Pt(30)
        x = mid_x + i * Pt(15)
        # 倒梯形（用圆角矩形代替）
        add_rect(slide, x, y, w, Pt(70), fill=color)
        add_textbox(slide, x + Pt(20), y + Pt(10), w - Pt(40), Pt(24),
                    f"{i+1}. {title}",
                    font_size=14, bold=True, color=theme.WHITE)
        add_textbox(slide, x + Pt(20), y + Pt(36), w - Pt(40), Pt(24),
                    desc,
                    font_size=11, color=theme.WHITE)

    # ---- 右栏：本项目定位（4 个差异化点）
    right_x = Pt(880)
    right_w = Pt(380)
    add_textbox(slide, right_x, col_top, right_w, Pt(30),
                "本项目定位：4 大差异化",
                font_size=15, bold=True, color="#FA8C16")
    diffs = [
        ("多智能体协同", "5 类智能体分工协作"),
        ("6 维动态画像", "随学随新、贴合个体"),
        ("结构化路径", "题库/模块双向同步"),
        ("流式可视化", "思考过程可折叠"),
    ]
    diff_top = col_top + Pt(50)
    for i, (title, desc) in enumerate(diffs):
        y = diff_top + i * Pt(60)
        # 左侧色块
        add_rect(slide, right_x, y + Pt(4), Pt(8), Pt(40), fill="#FA8C16")
        add_textbox(slide, right_x + Pt(20), y, right_w - Pt(20), Pt(24),
                    title,
                    font_size=14, bold=True, color=theme.PRIMARY_DARK)
        add_textbox(slide, right_x + Pt(20), y + Pt(24), right_w - Pt(20), Pt(24),
                    desc,
                    font_size=11, color=theme.TEXT_MUTED)