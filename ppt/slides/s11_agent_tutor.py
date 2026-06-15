"""
第 11 页 · 辅导答疑智能体（Tutor Agent）🟣
"""

import os
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from slides.s08_agent_profile import _screenshot


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=3, page_num=11)

    color = theme.AGENT_COLORS["tutor"]
    name = theme.AGENT_NAMES_CN["tutor"]
    en = theme.AGENT_NAMES_EN["tutor"]
    emoji = theme.AGENT_EMOJI["tutor"]

    add_color_block(slide, Pt(80), Pt(70), Pt(10), Pt(28), color)
    add_textbox(slide, left=Pt(100), top=Pt(70), width=Pt(700), height=Pt(38),
                text=f"{name} {emoji}", font_size=24, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, left=Pt(100), top=Pt(110), width=Pt(600), height=Pt(20),
                text=en, font_size=12, color=theme.TEXT_SUBTLE)

    _screenshot(slide, theme.SCREENSHOTS["tutor"], Pt(80), Pt(160), Pt(560), Pt(440))

    right_x = Pt(680)
    right_w = Pt(540)

    add_card(slide, right_x, Pt(160), right_w, Pt(110), fill=theme.ACCENT_BG, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(170), right_w - Pt(32), Pt(24),
                text="▶ 角色定位", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(196), right_w - Pt(32), Pt(70),
                text="一对一智能辅导，根据画像调整回答风格和深度；支持追问链。",
                font_size=12, color=theme.TEXT)

    add_card(slide, right_x, Pt(280), right_w, Pt(160), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(290), right_w - Pt(32), Pt(24),
                text="▶ 4 种解答模式", font_size=13, bold=True, color=color)
    modes = [("文字", "Markdown 渲染"), ("图解", "Mermaid/ASCII 流程"),
             ("视频", "脚本 + 时间戳"), ("代码", "可执行片段 + 注释")]
    for i, (cn_t, desc) in enumerate(modes):
        col = i % 2
        row = i // 2
        x = right_x + Pt(16) + col * Pt(260)
        y = Pt(320) + row * Pt(40)
        add_rect(slide, x, y + Pt(4), Pt(8), Pt(8), fill=color)
        add_textbox(slide, x + Pt(16), y, Pt(70), Pt(20),
                    text=cn_t, font_size=12, bold=True, color=color)
        add_textbox(slide, x + Pt(86), y, Pt(180), Pt(20),
                    text=desc, font_size=11, color=theme.TEXT)

    add_card(slide, right_x, Pt(450), right_w, Pt(150), fill=theme.WHITE, border=color, border_width=1.0)
    add_textbox(slide, right_x + Pt(16), Pt(460), right_w - Pt(32), Pt(24),
                text="▶ 工程优化（5 项）", font_size=13, bold=True, color=color)
    add_textbox(slide, right_x + Pt(16), Pt(488), right_w - Pt(32), Pt(110),
                text="• 画像注入 system prompt\n"
                     "• 缓存去重（问题+模式双键）\n"
                     "• 追问链 parentId/followUpIds\n"
                     "• 点踩重新生成（含原因分析）\n"
                     "• AbortSignal 取消未完成请求",
                font_size=12, color=theme.TEXT)
