"""
第 4 页 · 项目背景 · 学术商务版。
三栏布局 + 底部金句条（Y=494）。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title, apply_chrome_v2, add_bottom_bar
from components.shapes import add_card


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    apply_chrome_v2(slide, chapter_idx=1, page_num=4)

    add_page_title(slide, "项目背景与差异化定位",
                   subtitle="AI 教育市场趋势  ·  传统平台痛点分析  ·  本项目竞争优势")

    col_w = Pt(280)
    col_top = Pt(130)

    # ── 左栏：AI教育市场趋势 ──
    x1 = Pt(40)
    add_textbox(slide, x1, col_top, col_w, Pt(24),
                text="AI 教育市场趋势", font_size=17, bold=True, color=theme.NAVY)
    add_rect(slide, x1, col_top + Pt(26), Pt(28), Pt(2.5), fill=theme.NAVY)

    trends = [
        ("60%+", "学生认为AI个性化辅导\n优于传统课堂"),
        ("3 倍", "近2年自适应学习平台\n用户增长近3倍"),
        ("80%", "高校预计2026年引入\nAI教学辅助系统"),
    ]
    for i, (num, desc) in enumerate(trends):
        y = col_top + Pt(44) + i * Pt(78)
        add_card(slide, x1, y, col_w, Pt(68), fill=theme.LIGHT_GRAY)
        add_textbox(slide, x1 + Pt(14), y + Pt(6), Pt(70), Pt(30),
                    text=num, font_size=26, bold=True, color=theme.NAVY)
        add_textbox(slide, x1 + Pt(90), y + Pt(8), col_w - Pt(108), Pt(52),
                    text=desc, font_size=13, color=theme.DARK_TEXT)

    # 趋势结论
    add_card(slide, x1, col_top + Pt(286), col_w, Pt(42),
             fill=theme.NAVY, border=None)
    add_textbox(slide, x1 + Pt(10), col_top + Pt(290), col_w - Pt(20), Pt(34),
                text="AI+教育 2026 年进入\n高速增长期", font_size=14, bold=True,
                color=theme.WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ── 中栏：传统平台痛点 ──
    x2 = Pt(340)
    add_textbox(slide, x2, col_top, col_w, Pt(24),
                text="传统平台 3 大痛点", font_size=17, bold=True, color=theme.RED)

    pains = [
        ("资源繁杂分散", "题库/视频/文档分散各处，\n学生花40%+时间查找资源"),
        ("全班统一进度", "无个性化调整机制，\n70%+学生进度不匹配"),
        ("反馈严重滞后", "错题需等老师批改讲评，\n延迟使效率降低35%+"),
    ]
    for i, (title, desc) in enumerate(pains):
        y = col_top + Pt(44) + i * Pt(105)
        add_card(slide, x2, y, col_w, Pt(92), fill=theme.LIGHT_GRAY, border=theme.RED, border_width=0.75)
        add_textbox(slide, x2 + Pt(14), y + Pt(8), col_w - Pt(28), Pt(24),
                    text=f"0{i+1}  {title}", font_size=15, bold=True, color=theme.RED)
        add_textbox(slide, x2 + Pt(14), y + Pt(38), col_w - Pt(28), Pt(48),
                    text=desc, font_size=12, color=theme.DARK_TEXT)

    # ── 右栏：本项目4大差异化 ──
    x3 = Pt(640)
    add_textbox(slide, x3, col_top, col_w, Pt(24),
                text="本项目 4 大差异化", font_size=17, bold=True, color=theme.NAVY)

    diffs = [
        "多智能体协同：5类智能体分工协作\n事件总线完全解耦通信",
        "6维动态画像：对话式构建+随学随新\n画像注入所有下游智能体",
        "结构化路径：节点绑定题库/模块ID\n80%阈值自动标记+双向同步",
        "流式思考可视化：SSE打字机效果\n<thinking>块折叠+AbortSignal中断",
    ]
    for i, desc in enumerate(diffs):
        y = col_top + Pt(44) + i * Pt(80)
        add_card(slide, x3, y, col_w, Pt(70), fill=theme.LIGHT_GRAY)
        add_rect(slide, x3, y, Pt(4), Pt(70), fill=theme.NAVY)
        add_textbox(slide, x3 + Pt(14), y + Pt(6), Pt(28), Pt(24),
                    text=str(i + 1), font_size=18, bold=True, color=theme.NAVY)
        add_textbox(slide, x3 + Pt(44), y + Pt(6), col_w - Pt(62), Pt(58),
                    text=desc, font_size=12, color=theme.DARK_TEXT)

    # ─── 底部金句 ───
    add_bottom_bar(slide, "从\"统一供给\"到\"千人千面\"：打造每个学生专属的 AI 学习智能体",
                   highlight_words=["千人千面", "专属"])
