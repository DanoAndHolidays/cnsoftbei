"""
第 17 页 · 创新点小结。

5 个圆角胶囊（创新点编号 + 名称）+ 1 行简述。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block, add_capsule


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=5, page_num=17)

    add_page_title(slide, "创新点小结", subtitle="5 项创新点：多智能体框架 · 动态画像 · 结构化路径 · 流式思考 · 工程优化")

    innovations = [
        ("01", "5 智能体协同框架",     "MultiAgentScheduler 统一调度，事件总线解耦通信，6 worker 并行",       "#FA8C16"),
        ("02", "6 维动态画像",         "对话式构建 + 做题反馈回流，随学随新、跨页面共享",                    theme.PRIMARY),
        ("03", "结构化路径节点",       "StructuredLearningNode 绑定题库模块，80% 阈值自动标记",               "#52C41A"),
        ("04", "流式思考可视化",       "<thinking> 块折叠展示，typing 效果 + AbortSignal 取消",              "#722ED1"),
        ("05", "Tutor 5 项工程优化",   "画像注入 + 缓存去重 + 追问链 + 点踩重生 + 取消请求",                  "#13C2C2"),
    ]

    # 上方 5 个胶囊（编号 + 名称）
    cap_w = Pt(220)
    cap_h = Pt(60)
    cap_top = Pt(200)
    cap_gap = Pt(20)
    for i, (num, name, _desc, color) in enumerate(innovations):
        x = Pt(80) + i * (cap_w + cap_gap)
        # 编号小圆
        add_color_block(slide, x, cap_top, cap_w, Pt(8), color)
        add_card(slide, x, cap_top + Pt(8), cap_w, cap_h, fill=theme.WHITE, border=color, border_width=1.5)
        add_textbox(slide, x + Pt(12), cap_top + Pt(14), cap_w - Pt(24), Pt(24),
                    text=num, font_size=18, bold=True, color=color)
        add_textbox(slide, x + Pt(12), cap_top + Pt(36), cap_w - Pt(24), Pt(24),
                    text=name, font_size=12, bold=True, color=theme.PRIMARY_DARK)

    # 下方 5 个详细说明卡
    desc_top = Pt(310)
    desc_h = Pt(180)
    descs = [
        "统一 MultiAgentScheduler 调度 5 类智能体；通过事件总线解耦通信；\n6 worker 并行生成资源；失败自动重试 3 次；新增智能体只需 registerAgent。",
        "6 维画像（知识基础/认知风格/易错偏好/学习节奏/兴趣方向/学习习惯）；\n对话式构建；localStorage 持久化；做题反馈回流；下游智能体系统 prompt 注入。",
        "StructuredLearningNode 绑定 questionBankId + moduleId；\n12 条预定义路径 + AI 自由生成双轨；80% 完成度自动标记。",
        "SSE + ReadableStream 打字机效果；<thinking> 块可折叠；\nAbortSignal 支持中途取消；UI 实时显示 5 worker 状态。",
        "画像注入 + 缓存去重（问题+模式双键） + 追问链（parentId/followUpIds） +\n点踩重新生成（含原因分析） + AbortSignal 取消。",
    ]
    for i, (_num, _name, _desc, color) in enumerate(innovations):
        x = Pt(80) + i * (cap_w + cap_gap)
        add_card(slide, x, desc_top, cap_w, desc_h, fill=theme.ACCENT_BG, border=None)
        add_textbox(slide, x + Pt(12), desc_top + Pt(8), cap_w - Pt(24), desc_h - Pt(16),
                    text=descs[i], font_size=10, color=theme.TEXT)

    # 底部 callout
    add_card(slide, Pt(80), Pt(640), Pt(1130), Pt(40),
             fill=theme.PRIMARY_LIGHT, border="#13C2C2", border_width=1.0)
    add_textbox(slide, Pt(100), Pt(648), Pt(1090), Pt(24),
                text="5 项创新点覆盖架构层 / 数据层 / 交互层 / 工程层 4 个维度",
                font_size=13, bold=True, color=theme.PRIMARY_DARK, anchor=MSO_ANCHOR.MIDDLE)