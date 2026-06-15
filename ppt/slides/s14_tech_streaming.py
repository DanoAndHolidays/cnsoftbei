"""
第 14 页 · 关键技术 2：流式输出与思考过程可视化。
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_ANCHOR

from components import theme
from components.layout import add_textbox, add_rect, add_page_title
from components.shapes import add_card, add_color_block
from components.code_block import render_code_block
from slides.s08_agent_profile import _screenshot


def build(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    from components.layout import apply_chrome
    apply_chrome(slide, chapter_idx=4, page_num=14)

    add_color_block(slide, Pt(80), Pt(70), Pt(140), Pt(28), "#722ED1")
    add_textbox(slide, Pt(80), Pt(70), Pt(140), Pt(28),
                text="关键技术 02", font_size=14, bold=True, color=theme.WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left=Pt(232), top=Pt(70), width=Pt(900), height=Pt(38),
                text="流式输出与思考过程可视化",
                font_size=24, bold=True, color=theme.PRIMARY_DARK)
    add_textbox(slide, left=Pt(232), top=Pt(112), width=Pt(900), height=Pt(20),
                text="SSE + ReadableStream 打字机效果；思考过程可折叠；支持中途取消",
                font_size=12, color=theme.TEXT_MUTED)

    # 左：代码 + 流式截图
    code = """async function streamChatCompletion(
  messages: Message[],
  callbacks: { onChunk, onThinking },
  signal?: AbortSignal,
) {
  const res = await fetch('/anthropic/v1/messages', {
    method: 'POST',
    headers: { 'Content-Type': 'application/json' },
    body: JSON.stringify({ model, messages, stream: true }),
    signal,   // ← 关键：支持取消
  });
  const reader = res.body.getReader();
  const decoder = new TextDecoder();
  while (true) {
    const { done, value } = await reader.read();
    if (done) break;
    const text = decoder.decode(value);
    // 解析 SSE 事件，区分 thinking / content
    for (const event of parseSSE(text)) {
      if (event.type === 'thinking') callbacks.onThinking(event.delta);
      else if (event.type === 'content') callbacks.onChunk(event.delta);
    }
  }
}"""
    render_code_block(slide, Pt(80), Pt(180), Pt(600), Pt(300), code, font_size=10, lang_label="TS")

    # 流式截图（路径页的流式场景）
    _screenshot(slide, theme.SCREENSHOTS["path2"], Pt(80), Pt(500), Pt(600), Pt(180))

    # 右：3 个特性卡
    right_x = Pt(720)
    right_w = Pt(490)
    cards = [
        ("▶ 打字机效果", theme.PRIMARY, "逐 chunk 调用 setState，每帧渲染新字符。\n视觉上像 ChatGPT 一样逐字出现。"),
        ("▶ 思考过程可折叠", "#FA8C16", "AI 输出的 <thinking> 块折叠在\"💭 思考过程\" 标签下。\n默认折叠，用户点击展开看推理细节。"),
        ("▶ AbortSignal 取消", "#722ED1", "流式中点击取消 → 触发 controller.abort()。\nfetch 立即中断，UI 显示\"已取消\"状态。"),
    ]
    card_h = Pt(150)
    for i, (title, color, body) in enumerate(cards):
        y = Pt(180) + i * (card_h + Pt(20))
        add_card(slide, right_x, y, right_w, card_h, fill=theme.WHITE, border=color, border_width=1.5)
        add_color_block(slide, right_x, y, Pt(6), card_h, color)
        add_textbox(slide, right_x + Pt(16), y + Pt(12), right_w - Pt(32), Pt(28),
                    text=title, font_size=15, bold=True, color=color)
        add_textbox(slide, right_x + Pt(16), y + Pt(48), right_w - Pt(32), Pt(96),
                    text=body, font_size=12, color=theme.TEXT)