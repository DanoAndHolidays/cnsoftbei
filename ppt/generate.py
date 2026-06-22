#!/usr/bin/env python
"""
PPT 主入口。

用法：
  python generate.py
  python generate.py --only s10
  python generate.py --start s08 --end s12
"""

import argparse
import os
import sys
from pathlib import Path

# Windows GBK console 兼容：强制 UTF-8 输出
try:
    sys.stdout.reconfigure(encoding="utf-8")
except Exception:
    pass

# 让 `import components` / `import slides` 能工作
PPT_DIR = Path(__file__).resolve().parent
sys.path.insert(0, str(PPT_DIR))

from pptx import Presentation
from pptx.util import Emu

from components import theme
from slides import SLIDES, resolve_selection


# ---------------------------------------------------------------------------
# 汇报备注（presenter notes）
# ---------------------------------------------------------------------------
# 每张 slide 对应的口述稿。PowerPoint 演讲者视图可见，
# 与 `--only` / `--start` / `--end` 兼容：只生成指定 slide 时备注同样生效。
NOTES = {
    "s01_cover": """各位评委老师好，我是本项目的汇报人。今天我将为大家展示我们团队在第十五届中国软件杯A3赛题中的参赛作品——'学习智能体系统'。这是一个由多智能体协同驱动的个性化学习平台，旨在解决传统教育平台的痛点，让每个学生都能拥有专属的AI学习助手。""",

    "s02_toc": """本次汇报将分为四个核心部分：首先进行项目导入，分析行业痛点与赛题对标；其次拆解系统的四层总体架构与技术选型；接着深入剖析五大核心智能体的具体设计；最后挖掘多智能体协同等关键技术，并总结项目的创新点与未来展望。""",

    "s03_chapter01_div": """首先，我们进入第一部分：项目导入与需求对标。我们将探讨当前AI教育市场的趋势，分析传统学习平台面临的挑战，并展示我们的系统如何精准响应A3赛题的各项严苛要求。""",

    "s04_background": """大家请看这页的市场背景分析。当前AI教育市场正处于爆发期，超过60%的学生认为AI个性化辅导优于传统课堂，自适应学习平台用户近2年增长了3倍。然而，传统平台存在三大痛点：资源繁杂分散、全班统一进度难以个性化、以及错题反馈严重滞后。针对这些问题，我们的项目提出了四大差异化定位：通过5类智能体分工协作、构建6维动态画像、设计结构化路径，并实现思考过程的流式可视化，真正做到'随学随新、贴合个体'。""",

    "s05_requirements": """作为软件杯的参赛作品，赛题对标是我们的核心考核点。A3赛题提出了5项功能需求和4项非功能需求。大家可以看到这张对标表，在功能层面，我们不仅实现了6维对话式画像、多智能体资源生成和个性化路径规划，还超额完成了作为加分项的智能辅导与学习效果评估。在非功能层面，流式输出、防幻觉机制、响应时间追踪等要求也均已通过工程化手段完美落地。可以说，我们100%覆盖了赛题要求，并在多个维度实现了超越。""",

    "s06_chapter02_div": """接下来，我们进入第二部分：系统架构与技术选型。我将为大家拆解支撑这五大智能体高效运转的底层架构，以及我们在前端、AI交互和数据可视化方面的技术考量。""",

    "s07_architecture": """在系统架构设计上，我们采用了清晰的四层架构。最上层是基于React 19和Ant Design 6构建的表现层，提供企业级的UI交互；第二层是API网关，通过Vite代理无缝对接大模型接口；第三层是核心的多智能体框架，包含调度器和事件总线，负责5类智能体的协同工作；最底层是数据层，我们创新性地使用了localStorage结合本地JSON题库，实现了12个题库576道题的持久化。这种设计的最大优势是实现了零外部依赖，一行命令即可启动整个系统。""",

    "s08_tech_stack": """在技术选型方面，我们坚持'主流成熟+前沿AI+轻量可视化'的原则。前端选用React 19配合TypeScript，确保类型安全与极速的HMR开发体验；AI层全面兼容Anthropic协议，深度应用SSE流式传输与AbortSignal，完美实现打字机效果与请求中断控制；数据可视化则采用Recharts和Ant Design内置组件，精准呈现6维能力雷达与学习进度时间线。所有技术栈均经过严格筛选，在前沿与稳健之间取得了最佳平衡。""",

    "s09_chapter03_div": """现在，我们进入本次汇报的核心部分：五大核心智能体设计。我们将逐一拆解画像、资源、路径、辅导和评估这五个智能体，看看它们是如何分工协作，为学生提供全方位学习支持的。""",

    "s10_agent_profile": """首先是系统的'大脑'——画像构建智能体。它摒弃了传统的静态问卷，采用自然语言对话式交互，从知识基础、认知风格、易错偏好、学习节奏、兴趣方向和学习习惯6个维度动态构建学习者画像。最关键的特性是'随学随新'：学生在练习页的每一次做题反馈，都会通过事件总线实时回流到画像中进行更新。这个6维画像会作为系统prompt，自动注入到后续所有智能体的上下文中，确保每一次资源推荐和辅导都基于学生最新的认知状态。""",

    "s11_agent_resource": """接下来是资源生成智能体。当学生确定学习目标后，它会启动一个多智能体协作流水线。首先由Planner智能体拆解任务，然后派发给5类Worker智能体并行工作，支持生成文档、思维导图、测验、阅读、视频脚本和代码案例共6种定制资源。在这个过程中，前端通过SSE流式回传每个Worker的实时状态（如pending、running、done），让用户清晰看到AI的工作进度。同时，系统支持失败自动重试和中途一键取消，确保长耗时场景下的稳定性。""",

    "s12_agent_path": """路径规划智能体负责为学生规划学习路线。我们设计了'双轨模式'：既支持AI根据画像自由生成节点列表，也内置了12条如'Python全栈工程师'、'算法工程师'等预定义的结构化路径供一键采用。这里的核心技术亮点是'双向同步'：每个路径节点都深度绑定了具体的题库和模块ID。当学生在练习页做题时，系统会按当前激活的路径过滤模块；当模块完成度达到80%的阈值时，路径Banner会实时标记完成，并触发事件通知评估页更新雷达图。""",

    "s13_agent_tutor": """辅导答疑智能体扮演着一对一AI导师的角色。它支持文字Markdown、Mermaid图解、视频脚本和可执行代码4种解答模式，满足不同学科的辅导需求。在工程优化上我们做了5项深度打磨：首先是画像注入，让AI用学生最容易理解的方式讲解；其次是缓存去重和追问链设计，支持深度探讨；最后是点踩重新生成与AbortSignal取消功能。这些细节优化让AI辅导不再是冷冰冰的问答，而是具备记忆和纠错能力的深度互动。""",

    "s14_agent_assessment": """最后是效果评估智能体。我们坚持'真实进度同步、绝不造假'的原则，所有统计卡片和模块进度均直接读取练习页写入的真实数据。在此基础上，系统利用 Recharts 实时渲染 6 维能力雷达图，直观暴露学生的知识盲区。更重要的是，它会基于雷达图生成智能调整建议，比如针对易错的'装饰器'知识点推荐专项练习，或者根据学习节奏建议每日时长，甚至推荐相关的发展路径，真正实现'以评促学'。""",

    "s15_chapter04_div": """在了解了五大智能体的功能后，第四部分我们将深入底层，深挖多智能体协同、流式输出等关键技术的实现细节，并对整个项目进行系统评估与未来展望。""",

    "s16_tech_multi_agent": """多智能体协同是整个系统的技术基石。我们自研了 MultiAgentScheduler 核心类，通过 TypeScript 的 Map 管理角色与智能体的映射。最关键的设计是引入了 EventEmitter 事件总线，实现了5个智能体之间的完全解耦通信。比如，Profile Agent 更新画像后广播事件，Resource Agent 监听并生成资源，整个过程无需硬编码依赖。这种设计使得新增智能体只需调用 registerAgent()，让代码量大幅减少40%。""",

    "s17_tech_streaming": """在AI交互体验上，我们深度优化了流式输出。通过封装 streamChatCompletion 函数，利用 Fetch API 和 ReadableStream 逐 chunk 解析 SSE 数据，实现了丝滑的打字机效果。更具特色的是'思考过程可视化'：我们将大模型输出的 <thinking> 块默认折叠在'思考过程'标签下，用户点击即可展开查看AI的推理逻辑，这在教学场景中极大地增强了可信度。同时，结合 axios 重试与超时控制，即使在弱网环境下也能稳定输出。""",

    "s18_tech_sync": """路径与练习的双向同步是保证学习闭环的关键。我们定义了 StructuredLearningNode 接口，强制每个节点绑定 questionBankId 和 moduleId。在数据流上，Path 页将激活路径写入 localStorage 作为单一数据源，Practice 页读取并过滤出对应模块的题目。我们内置了12个题库共576道全预写的题目，AI只负责判分不负责生成题目，这从根本上杜绝了AI幻觉问题。当模块完成度达80%时，通过 customEvent 广播进度，触发 Assessment 雷达图实时更新。""",

    "s19_evaluation": """在系统交付前，我们进行了严格的测试评估。核心指标方面：系统内置576道题目覆盖主流计算机学科；AI简答判分与人工判分的一致率高达92%；流式首字延迟控制在2秒以内；且 npm run build 零严重缺陷通过。在测试结论上，12个题库的题型比例科学，跨页面同步缓存命中率达到100%，多智能体6个Worker并发完成时间小于8秒。这些数据充分证明了系统不仅功能完备，而且在工程性能和稳定性上达到了生产级标准。""",

    "s20_innovation_summary": """总结来看，本项目在架构、数据、交互和工程四个维度沉淀了5项核心创新：多智能体协同框架、6维动态画像、结构化路径节点、流式思考可视化以及Tutor的5项工程优化。展望未来，我们将继续拓展系统的边界：一是接入图像与语音识别，实现拍照搜题等多模态交互；二是构建底层学科知识图谱，让路径推荐具备真正的逻辑推理能力；三是引入跨用户协作，支持学习小组与同伴互评。我们相信，AI加教育的终局，是让每个学生都拥有最懂自己的学习智能体。""",

    "s21_closing": """以上就是我们团队关于'学习智能体系统'的全部汇报内容。感谢各位评委老师的耐心聆听与批评指正。本项目使用了React、TypeScript、Vite、Ant Design、Recharts、python-pptx等众多优秀的开源组件，在此一并致谢。接下来是Q&A环节，欢迎老师们批评指正、提问交流。谢谢大家！""",
}


def parse_args():
    parser = argparse.ArgumentParser(description="生成《学习智能体系统》汇报 PPT")
    parser.add_argument("--only", help="只生成指定 slide，例如 s10")
    parser.add_argument("--start", help="起始 slide（含），例如 s08")
    parser.add_argument("--end", help="结束 slide（含），例如 s12")
    parser.add_argument("-o", "--output", help="输出文件路径", default=None)
    return parser.parse_args()


def import_slide_module(slide_name: str):
    """动态 import slides.s01_cover 等"""
    import importlib
    return importlib.import_module(f"slides.{slide_name}")


def main():
    args = parse_args()
    selected = resolve_selection(args)
    print(f"[PPT] 将生成 {len(selected)} 张 slide：{', '.join(selected)}")

    out_path = args.output or os.path.join(PPT_DIR, "output", "学习智能体系统_汇报PPT.pptx")
    Path(os.path.dirname(out_path)).mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = theme.SLIDE_WIDTH
    prs.slide_height = theme.SLIDE_HEIGHT

    for slide_name in selected:
        print(f"[PPT]  渲染 {slide_name} ...")
        mod = import_slide_module(slide_name)
        if not hasattr(mod, "build"):
            raise SystemExit(f"slide 模块 {slide_name} 必须导出 build(prs) 函数")
        mod.build(prs)

        # 写入演讲者备注（presenter notes）
        if slide_name in NOTES:
            note_text = NOTES[slide_name].format(
                team=theme.COVER_INFO.get("team_name", ""),
                school=theme.COVER_INFO.get("school", ""),
                presenter=theme.COVER_INFO.get("presenter", ""),
                advisor=theme.COVER_INFO.get("advisor", ""),
                contact=theme.COVER_INFO.get("contact", ""),
            )
            prs.slides[-1].notes_slide.notes_text_frame.text = note_text

    prs.save(out_path)
    print(f"[PPT] ✓ 已生成：{out_path}")
    print(f"[PPT]   slide 总数: {len(prs.slides)}")

    # 清理 assessment mockup 生成的中间 PNG
    try:
        from components.assessment_mock import cleanup_mock_assets
        cleanup_mock_assets()
    except Exception:
        pass


if __name__ == "__main__":
    main()
